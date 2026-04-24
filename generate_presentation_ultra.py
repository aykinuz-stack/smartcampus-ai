#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SmartCampus AI — ULTRA PREMIUM Investor/Sales Presentation Generator
=====================================================================
World-class 50+ slide 16:9 PowerPoint for school owners & investors.
Dark navy + gold design language. Fully self-contained — no image files needed.

Output: SmartCampusAI_Sunum_ULTRA.pptx
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy

# ====================================================================
# DESIGN TOKENS
# ====================================================================
# Backgrounds
DARK_NAVY    = RGBColor(0x02, 0x06, 0x17)   # #020617  (deepest)
NAVY         = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A  (primary bg)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B  (card fill)
CARD_BG_ALT  = RGBColor(0x15, 0x20, 0x33)   # #152033  (darker card)

# Accents
GOLD         = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B  (primary accent)
LIGHT_GOLD   = RGBColor(0xFB, 0xBF, 0x24)   # #FBBF24
DARK_GOLD    = RGBColor(0xD9, 0x77, 0x06)   # #D97706

# Text
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0
LIGHT_GRAY   = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8  (body text)
MED_GRAY     = RGBColor(0x64, 0x74, 0x8B)   # #64748B
DARK_GRAY    = RGBColor(0x33, 0x3D, 0x4D)   # #333D4D

# Semantic colors
BLUE         = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6
CYAN         = RGBColor(0x06, 0xB6, 0xD4)   # #06B6D4
GREEN        = RGBColor(0x10, 0xB9, 0x81)   # #10B981
EMERALD      = RGBColor(0x34, 0xD3, 0x99)   # #34D399
PURPLE       = RGBColor(0x8B, 0x5C, 0xF6)   # #8B5CF6
RED          = RGBColor(0xEF, 0x44, 0x44)   # #EF4444
PINK         = RGBColor(0xEC, 0x48, 0x99)   # #EC4899
ORANGE       = RGBColor(0xF9, 0x73, 0x16)   # #F97316
TEAL         = RGBColor(0x14, 0xB8, 0xA6)   # #14B8A6

# Dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins
MARGIN_LEFT   = Inches(0.8)
MARGIN_RIGHT  = Inches(0.8)
MARGIN_TOP    = Inches(0.6)
CONTENT_WIDTH = Inches(11.733)  # 13.333 - 0.8*2

# Footer
FOOTER_TEXT = "SmartCampus AI  |  Gizli & Ozel  |  2026"
FOOTER_Y    = Inches(7.05)

# ====================================================================
# PRESENTATION INIT
# ====================================================================
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ====================================================================
# HELPER FUNCTIONS
# ====================================================================

def _rgb_str(c: RGBColor) -> str:
    """Return '0F172A' style hex string."""
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def add_solid_bg(slide, color=NAVY):
    """Solid color background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, c1=DARK_NAVY, c2=NAVY):
    """Two-stop linear gradient background (top -> bottom)."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0


def new_slide():
    """Create a blank slide with gradient bg, gold separator line, and footer."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_gradient_bg(slide)
    # Gold separator line near bottom
    _add_line(slide, Inches(0.6), Inches(6.85), Inches(12.133), Inches(6.85),
              color=GOLD, width=Pt(0.75))
    # Footer text
    _add_text_box(slide, Inches(0.8), FOOTER_Y, Inches(11.733), Inches(0.35),
                  FOOTER_TEXT, Pt(8), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                  font_name="Segoe UI")
    return slide


def _add_line(slide, x1, y1, x2, y2, color=GOLD, width=Pt(1)):
    """Add a straight line connector."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width


def _add_shape(slide, shape_type, left, top, width, height,
               fill_color=None, border_color=None, border_width=Pt(1)):
    """Add an auto shape with optional fill and border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height,
                      fill_color=CARD_BG, border_color=None,
                      border_width=Pt(1), radius=Inches(0.15)):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05  # ~5% corner radius
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                  color=WHITE, alignment=PP_ALIGN.LEFT, bold=False,
                  font_name="Segoe UI", line_spacing=1.2, space_after=Pt(0)):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    if line_spacing != 1.2:
        p.line_spacing = line_spacing
    return txBox


def _add_rich_text_box(slide, left, top, width, height, runs,
                       alignment=PP_ALIGN.LEFT, line_spacing=1.15):
    """Add text box with multiple styled runs.
    runs = [(text, font_size, color, bold, font_name), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    for i, run_def in enumerate(runs):
        text, size, color, bold = run_def[0], run_def[1], run_def[2], run_def[3]
        fname = run_def[4] if len(run_def) > 4 else "Segoe UI"
        if i == 0:
            r = p.runs[0] if p.runs else p.add_run()
            r.text = text
        else:
            r = p.add_run()
            r.text = text
        r.font.size = size
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = fname
    return txBox


def _add_multiline_text(slide, left, top, width, height, lines,
                        alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """Add text box with multiple paragraphs.
    lines = [(text, font_size, color, bold, font_name), ...]
    Each tuple becomes its own paragraph.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_def in enumerate(lines):
        text, size, color, bold = line_def[0], line_def[1], line_def[2], line_def[3]
        fname = line_def[4] if len(line_def) > 4 else "Segoe UI"
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = fname
        p.alignment = alignment
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
    return txBox


def _add_bullet_list(slide, left, top, width, height, items,
                     font_size=Pt(14), color=LIGHT_GRAY, bullet_color=GOLD,
                     alignment=PP_ALIGN.LEFT, spacing=1.4, bold=False):
    """Add a bullet list. items = list of strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Gold bullet char + text
        r1 = p.add_run()
        r1.text = "\u25C6  "
        r1.font.size = Pt(10)
        r1.font.color.rgb = bullet_color
        r1.font.name = "Segoe UI"
        r2 = p.add_run()
        r2.text = item
        r2.font.size = font_size
        r2.font.color.rgb = color
        r2.font.bold = bold
        r2.font.name = "Segoe UI"
        p.alignment = alignment
        p.line_spacing = spacing
        p.space_after = Pt(3)
    return txBox


def _add_icon_bullet_list(slide, left, top, width, height, items,
                          font_size=Pt(13), color=LIGHT_GRAY,
                          spacing=1.4):
    """Bullet list with emoji icons. items = [(emoji, text), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (icon, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{icon}  "
        r1.font.size = Pt(16)
        r1.font.name = "Segoe UI Emoji"
        r2 = p.add_run()
        r2.text = text
        r2.font.size = font_size
        r2.font.color.rgb = color
        r2.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_after = Pt(2)
    return txBox


def _add_stat_card(slide, left, top, width, height, number, label,
                   num_color=GOLD, label_color=LIGHT_GRAY,
                   bg_color=CARD_BG, border_color=None):
    """Big number stat card."""
    card = _add_rounded_rect(slide, left, top, width, height,
                             fill_color=bg_color, border_color=border_color)
    # Number
    _add_text_box(slide, left, top + Inches(0.15), width, Inches(0.55),
                  str(number), Pt(36), num_color, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Black")
    # Label
    _add_text_box(slide, left, top + Inches(0.7), width, Inches(0.4),
                  label, Pt(11), label_color, PP_ALIGN.CENTER, bold=False)
    return card


def _add_small_stat(slide, left, top, width, height, number, label,
                    num_color=GOLD, label_color=LIGHT_GRAY, bg_color=CARD_BG):
    """Smaller stat card."""
    card = _add_rounded_rect(slide, left, top, width, height,
                             fill_color=bg_color, border_color=None)
    _add_text_box(slide, left, top + Inches(0.08), width, Inches(0.4),
                  str(number), Pt(24), num_color, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Black")
    _add_text_box(slide, left, top + Inches(0.45), width, Inches(0.3),
                  label, Pt(9), label_color, PP_ALIGN.CENTER, bold=False)
    return card


def _slide_title(slide, title, subtitle=None, title_y=Inches(0.45),
                 underline=True, title_size=Pt(32)):
    """Add slide title with optional gold underline and subtitle."""
    _add_text_box(slide, MARGIN_LEFT, title_y, CONTENT_WIDTH, Inches(0.55),
                  title, title_size, WHITE, PP_ALIGN.LEFT, bold=True,
                  font_name="Segoe UI Black")
    if underline:
        _add_line(slide, MARGIN_LEFT, title_y + Inches(0.6),
                  Inches(4.5), title_y + Inches(0.6),
                  color=GOLD, width=Pt(3))
    if subtitle:
        _add_text_box(slide, MARGIN_LEFT, title_y + Inches(0.7),
                      CONTENT_WIDTH, Inches(0.4),
                      subtitle, Pt(14), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)


def _section_label(slide, left, top, text, color=GOLD, size=Pt(10)):
    """Small uppercase section label."""
    _add_text_box(slide, left, top, Inches(5), Inches(0.25),
                  text.upper(), size, color, PP_ALIGN.LEFT, bold=True,
                  font_name="Segoe UI Semibold")


def _feature_card(slide, left, top, width, height, emoji, title, desc,
                  title_color=WHITE, desc_color=LIGHT_GRAY,
                  bg_color=CARD_BG, border_color=None):
    """Feature card with emoji, title, description."""
    _add_rounded_rect(slide, left, top, width, height,
                      fill_color=bg_color, border_color=border_color)
    # Emoji
    _add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                  Inches(0.5), Inches(0.4),
                  emoji, Pt(22), WHITE, PP_ALIGN.LEFT, bold=False,
                  font_name="Segoe UI Emoji")
    # Title
    _add_text_box(slide, left + Inches(0.2), top + Inches(0.55),
                  width - Inches(0.4), Inches(0.35),
                  title, Pt(13), title_color, PP_ALIGN.LEFT, bold=True)
    # Description
    _add_text_box(slide, left + Inches(0.2), top + Inches(0.9),
                  width - Inches(0.4), height - Inches(1.05),
                  desc, Pt(10), desc_color, PP_ALIGN.LEFT, bold=False)


def _role_card(slide, left, top, width, height, emoji, role_name,
               features, accent_color=GOLD):
    """Role card with emoji header and feature list."""
    _add_rounded_rect(slide, left, top, width, height,
                      fill_color=CARD_BG, border_color=accent_color,
                      border_width=Pt(1.5))
    # Accent bar at top
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               left + Inches(0.05), top + Inches(0.05),
               width - Inches(0.1), Inches(0.06),
               fill_color=accent_color)
    # Emoji + Role name
    _add_rich_text_box(
        slide, left + Inches(0.15), top + Inches(0.2),
        width - Inches(0.3), Inches(0.5),
        [
            (emoji + "  ", Pt(20), WHITE, False, "Segoe UI Emoji"),
            (role_name, Pt(15), WHITE, True, "Segoe UI Black"),
        ],
        alignment=PP_ALIGN.CENTER
    )
    # Features
    _add_icon_bullet_list(
        slide, left + Inches(0.15), top + Inches(0.75),
        width - Inches(0.3), height - Inches(0.9),
        [("\u2022", f) for f in features],
        font_size=Pt(10), color=LIGHT_GRAY, spacing=1.25
    )


def _module_slide(slide, emoji, title, subtitle, stats, features,
                  highlight_text=None):
    """Standard module detail slide layout."""
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "MODUL DETAY")
    # Title with emoji
    _add_rich_text_box(
        slide, MARGIN_LEFT, Inches(0.55), CONTENT_WIDTH, Inches(0.6),
        [
            (emoji + "  ", Pt(30), GOLD, False, "Segoe UI Emoji"),
            (title, Pt(30), WHITE, True, "Segoe UI Black"),
        ]
    )
    # Gold underline
    _add_line(slide, MARGIN_LEFT, Inches(1.2),
              Inches(5), Inches(1.2), color=GOLD, width=Pt(3))
    # Subtitle
    if subtitle:
        _add_text_box(slide, MARGIN_LEFT, Inches(1.35),
                      Inches(7), Inches(0.35),
                      subtitle, Pt(13), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)

    # Stats row
    if stats:
        stat_w = Inches(1.8)
        stat_h = Inches(1.0)
        stat_gap = Inches(0.2)
        sx = MARGIN_LEFT
        sy = Inches(1.85)
        for num, label in stats:
            _add_stat_card(slide, sx, sy, stat_w, stat_h, num, label,
                           border_color=DARK_GOLD)
            sx += stat_w + stat_gap

    # Features (right side or below stats)
    feat_y = Inches(3.1) if stats else Inches(1.85)
    if features:
        _add_icon_bullet_list(
            slide, MARGIN_LEFT, feat_y, Inches(5.5), Inches(3.5),
            features, font_size=Pt(12), color=OFF_WHITE, spacing=1.35
        )

    # Highlight box (right side)
    if highlight_text:
        _add_rounded_rect(slide, Inches(7.5), Inches(1.85),
                          Inches(5.0), Inches(4.5),
                          fill_color=CARD_BG_ALT, border_color=GOLD,
                          border_width=Pt(1))
        _add_text_box(slide, Inches(7.7), Inches(2.0),
                      Inches(4.6), Inches(4.2),
                      highlight_text, Pt(11), LIGHT_GRAY, PP_ALIGN.LEFT,
                      bold=False, line_spacing=1.5)


def _comparison_row(slide, left, top, width, row_h, cols, is_header=False):
    """Draw a comparison table row."""
    col_w = width / len(cols)
    for i, (text, color) in enumerate(cols):
        x = left + col_w * i
        bg = CARD_BG_ALT if is_header else (CARD_BG if (top // row_h) % 2 == 0 else CARD_BG_ALT)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, x, top, col_w, row_h,
                   fill_color=bg if not is_header else DARK_GOLD,
                   border_color=DARK_GRAY, border_width=Pt(0.5))
        sz = Pt(10) if is_header else Pt(9)
        _add_text_box(slide, x + Inches(0.05), top, col_w - Inches(0.1), row_h,
                      text, sz, color, PP_ALIGN.CENTER,
                      bold=is_header)


# ====================================================================
# SLIDE 1: KAPAK
# ====================================================================
def slide_01_cover():
    slide = new_slide()
    add_gradient_bg(slide, DARK_NAVY, NAVY)

    # Top-left gold accent line
    _add_line(slide, Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.5),
              color=GOLD, width=Pt(4))
    _add_line(slide, Inches(0.5), Inches(0.5), Inches(0.5), Inches(2.5),
              color=GOLD, width=Pt(4))

    # Confidential badge
    badge = _add_rounded_rect(slide, Inches(9.8), Inches(0.4),
                              Inches(3.0), Inches(0.4),
                              fill_color=DARK_GOLD, border_color=GOLD)
    _add_text_box(slide, Inches(9.8), Inches(0.4), Inches(3.0), Inches(0.4),
                  "\U0001F512  GIZLI & OZEL DOKUMAN", Pt(10), WHITE,
                  PP_ALIGN.CENTER, bold=True)

    # Main title
    _add_text_box(slide, Inches(1.0), Inches(2.0), Inches(11.0), Inches(1.0),
                  "SmartCampus AI", Pt(60), WHITE, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Black")

    # Gold divider
    _add_line(slide, Inches(4.5), Inches(3.2), Inches(8.833), Inches(3.2),
              color=GOLD, width=Pt(4))

    # Tagline
    _add_text_box(slide, Inches(1.0), Inches(3.5), Inches(11.0), Inches(0.6),
                  "Egitimin Gelecegi", Pt(28), GOLD, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Light")

    # Description
    _add_text_box(slide, Inches(2.0), Inches(4.3), Inches(9.0), Inches(0.8),
                  "Turkiye'nin en kapsamli, yapay zeka destekli egitim yonetim platformu",
                  Pt(16), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)

    # Stats preview bar
    bar_y = Inches(5.3)
    bar_items = [("30+", "Web Modul"), ("99", "Mobil Sayfa"), ("143", "API Route"), ("5", "Kullanici Rolu")]
    bar_w = Inches(2.5)
    bar_gap = Inches(0.3)
    start_x = Inches(1.5)
    for i, (num, label) in enumerate(bar_items):
        x = start_x + i * (bar_w + bar_gap)
        _add_small_stat(slide, x, bar_y, bar_w, Inches(0.8),
                        num, label, num_color=GOLD, bg_color=CARD_BG_ALT)

    # Date & version
    _add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.0), Inches(0.35),
                  "Nisan 2026  |  v3.0  |  Yatirimci Sunumu", Pt(11),
                  MED_GRAY, PP_ALIGN.CENTER, bold=False)

    # Bottom-right gold accent
    _add_line(slide, Inches(9.833), Inches(6.6), Inches(12.833), Inches(6.6),
              color=GOLD, width=Pt(4))
    _add_line(slide, Inches(12.833), Inches(4.6), Inches(12.833), Inches(6.6),
              color=GOLD, width=Pt(4))


# ====================================================================
# SLIDE 2: NEDEN SmartCampus?
# ====================================================================
def slide_02_neden():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "PROBLEM TESPITI")
    _slide_title(slide, "Neden SmartCampus AI?",
                 "Ozel okullarin yasadigi 5 kritik problem", title_y=Inches(0.5))

    problems = [
        ("\U0001F4C9", "Verimsizlik & Zaman Kaybi",
         "Ogretmenler gunde 2+ saat idari isle ugrasir. Not girisi, yoklama, raporlama hep manuel.",
         RED),
        ("\U0001F4C4", "Kagit Isler & Maliyet",
         "Sinav basimi, form doldurmasi, arsivleme. Yillik binlerce sayfa kagit israfi.",
         ORANGE),
        ("\U0001F50C", "Iletisim Kopuklugu",
         "Veli-okul iletisimi tek yonlu. Veliler bilgiye gec ulasir. SMS maliyeti yuksek.",
         PURPLE),
        ("\U0001F4CA", "Takip Eksikligi",
         "Ogrenci basarisi, devamsizlik, risk analizleri yapilamaz. Veri daginiik.",
         BLUE),
        ("\u26A1", "Dijitallesme Baskisi",
         "Rakip kurumlar dijitallesirken geride kalmak kayit kaybina neden olur.",
         GOLD),
    ]

    card_w = Inches(2.15)
    card_h = Inches(3.5)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    start_y = Inches(2.5)

    for i, (emoji, title, desc, accent) in enumerate(problems):
        x = start_x + i * (card_w + gap)
        _add_rounded_rect(slide, x, start_y, card_w, card_h,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1.5))
        # Accent top bar
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   x + Inches(0.05), start_y + Inches(0.05),
                   card_w - Inches(0.1), Inches(0.06),
                   fill_color=accent)
        # Number circle
        _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.2),
                      Inches(0.4), Inches(0.4),
                      emoji, Pt(24), WHITE, PP_ALIGN.LEFT, bold=False,
                      font_name="Segoe UI Emoji")
        # Title
        _add_text_box(slide, x + Inches(0.15), start_y + Inches(0.7),
                      card_w - Inches(0.3), Inches(0.5),
                      title, Pt(12), WHITE, PP_ALIGN.LEFT, bold=True)
        # Desc
        _add_text_box(slide, x + Inches(0.15), start_y + Inches(1.3),
                      card_w - Inches(0.3), Inches(2.0),
                      desc, Pt(10), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)


# ====================================================================
# SLIDE 3: COZUM OZETI
# ====================================================================
def slide_03_cozum():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "COZUM")
    _slide_title(slide, "Tek Platform, Sinirsiz Guc",
                 title_y=Inches(0.5))

    # Value prop
    _add_text_box(slide, MARGIN_LEFT, Inches(1.3), Inches(11.0), Inches(0.5),
                  "Web + Mobil + API = Tam Entegre Egitim Ekosistemi",
                  Pt(18), GOLD, PP_ALIGN.LEFT, bold=True)

    # Three pillars
    pillars = [
        ("\U0001F4BB", "WEB PLATFORMU", "Streamlit",
         "30+ modul, 390+ sekme\nYonetici & ogretmen paneli\nAnalitik dashboard",
         BLUE),
        ("\U0001F4F1", "MOBIL UYGULAMA", "Flutter Android",
         "99 sayfa, 5 rol\nOffline mod, bildirimler\nUltra premium tasarim",
         GREEN),
        ("\u2699\uFE0F", "BACKEND API", "FastAPI",
         "143 route, JWT guvenlik\nMulti-tenant mimari\nSMS/Email entegrasyon",
         PURPLE),
    ]

    pw = Inches(3.6)
    ph = Inches(3.5)
    gap = Inches(0.35)
    sx = Inches(0.9)
    sy = Inches(2.2)

    for i, (emoji, title, tech, desc, accent) in enumerate(pillars):
        x = sx + i * (pw + gap)
        _add_rounded_rect(slide, x, sy, pw, ph,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(2))
        # Emoji
        _add_text_box(slide, x, sy + Inches(0.25), pw, Inches(0.5),
                      emoji, Pt(36), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        # Title
        _add_text_box(slide, x, sy + Inches(0.85), pw, Inches(0.35),
                      title, Pt(16), WHITE, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Tech badge
        badge_w = Inches(1.5)
        _add_rounded_rect(slide, x + (pw - badge_w) / 2, sy + Inches(1.3),
                          badge_w, Inches(0.3),
                          fill_color=accent, border_color=None)
        _add_text_box(slide, x + (pw - badge_w) / 2, sy + Inches(1.3),
                      badge_w, Inches(0.3),
                      tech, Pt(9), WHITE, PP_ALIGN.CENTER, bold=True)
        # Description
        _add_text_box(slide, x + Inches(0.3), sy + Inches(1.8),
                      pw - Inches(0.6), Inches(1.5),
                      desc, Pt(12), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                      line_spacing=1.5)

    # Bottom connector line
    _add_line(slide, sx + pw / 2, sy + ph, sx + pw / 2, sy + ph + Inches(0.3),
              color=GOLD, width=Pt(2))
    _add_line(slide, sx + pw + gap + pw / 2, sy + ph,
              sx + pw + gap + pw / 2, sy + ph + Inches(0.3),
              color=GOLD, width=Pt(2))
    _add_line(slide, sx + 2 * (pw + gap) + pw / 2, sy + ph,
              sx + 2 * (pw + gap) + pw / 2, sy + ph + Inches(0.3),
              color=GOLD, width=Pt(2))
    _add_line(slide, sx + pw / 2, sy + ph + Inches(0.3),
              sx + 2 * (pw + gap) + pw / 2, sy + ph + Inches(0.3),
              color=GOLD, width=Pt(2))


# ====================================================================
# SLIDE 4: RAKAMLARLA SmartCampus
# ====================================================================
def slide_04_numbers():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BUYUK RESIM")
    _slide_title(slide, "Rakamlarla SmartCampus AI", title_y=Inches(0.5))

    # Big stats grid (2 rows x 4 cols)
    stats = [
        ("30+", "Web Modul", GOLD),
        ("99", "Mobil Sayfa", GOLD),
        ("143", "API Route", GOLD),
        ("390+", "Toplam Sekme", GOLD),
        ("5", "Kullanici Rolu", LIGHT_GOLD),
        ("687+", "Veri Dosyasi", LIGHT_GOLD),
        ("8", "Sidebar Grubu", LIGHT_GOLD),
        ("7", "Sertifika Sablonu", LIGHT_GOLD),
    ]

    card_w = Inches(2.6)
    card_h = Inches(1.6)
    gap_x = Inches(0.25)
    gap_y = Inches(0.3)
    start_x = Inches(0.75)
    start_y = Inches(1.8)

    for i, (num, label, color) in enumerate(stats):
        row = i // 4
        col = i % 4
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        _add_rounded_rect(slide, x, y, card_w, card_h,
                          fill_color=CARD_BG, border_color=DARK_GOLD,
                          border_width=Pt(1))
        _add_text_box(slide, x, y + Inches(0.2), card_w, Inches(0.7),
                      num, Pt(42), color, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        _add_text_box(slide, x, y + Inches(0.95), card_w, Inches(0.35),
                      label, Pt(13), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)

    # Bottom quote
    _add_text_box(slide, Inches(1.0), Inches(5.8), Inches(11.0), Inches(0.5),
                  "\"Turkiye'nin en kapsamli egitim yonetim sistemi\"",
                  Pt(16), GOLD, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Light")
    _add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.0), Inches(0.3),
                  "Web + Mobil + API  |  Yapay Zeka Destekli  |  Multi-Tenant  |  Enterprise Kalite",
                  Pt(11), MED_GRAY, PP_ALIGN.CENTER, bold=False)


# ====================================================================
# SLIDE 5: ROLLER
# ====================================================================
def slide_05_roles():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "KULLANICI ROLLERI")
    _slide_title(slide, "5 Rol, Tek Platform", title_y=Inches(0.5))

    roles = [
        ("\U0001F451", "YONETICI", [
            "Tum modullere erisim",
            "Dashboard & analitik",
            "Butce & odeme takip",
            "Personel yonetimi",
            "Kurumsal raporlar",
        ], GOLD),
        ("\U0001F4DA", "OGRETMEN", [
            "Not & yoklama girisi",
            "Ders defteri & plan",
            "Sinav olusturma",
            "Odev yonetimi",
            "AI soru uretimi",
        ], BLUE),
        ("\U0001F9E0", "REHBER", [
            "Vaka yonetimi",
            "Mood check-in",
            "BEP takibi",
            "Gorusme kaydi",
            "Risk analizi",
        ], PURPLE),
        ("\U0001F468\u200D\U0001F469\u200D\U0001F467", "VELI", [
            "Cocuk takip paneli",
            "Randevu alma",
            "Odeme durumu",
            "Servis GPS takip",
            "Mesajlasma",
        ], GREEN),
        ("\U0001F393", "OGRENCI", [
            "Dashboard & notlar",
            "Online sinav",
            "AI asistan",
            "Dijital ogrenme",
            "Etkinlik & kulupler",
        ], CYAN),
    ]

    rw = Inches(2.2)
    rh = Inches(3.8)
    gap = Inches(0.2)
    sx = Inches(0.55)
    sy = Inches(1.8)

    for i, (emoji, name, features, accent) in enumerate(roles):
        x = sx + i * (rw + gap)
        _role_card(slide, x, sy, rw, rh, emoji, name, features, accent)

    # Bottom note
    _add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11.0), Inches(0.35),
                  "Her rol sadece yetkili oldugu modulleri gorur  |  Rol bazli erisim kontrolu",
                  Pt(11), MED_GRAY, PP_ALIGN.CENTER, bold=False)


# ====================================================================
# MODULE SLIDES (6–25)
# ====================================================================

def slide_06_kayit():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4CB", "Kayit Modulu", "Ogrenci kayit sureci basindan sonuna",
        [("6", "Grup"), ("34", "Sekme"), ("5", "Pipeline Adim"), ("3", "Kampanya Turu")],
        [
            ("\U0001F4CB", "Tam pipeline: Basvuru -> Gorusme -> Sinav -> Kabul -> Kayit -> Kesin Kayit"),
            ("\U0001F3AF", "Kampanya yonetimi: erken kayit, kardes, burslu indirimleri"),
            ("\U0001F916", "AI tabanli lead skorlama ve takip otomasyonu"),
            ("\U0001F4CA", "Donem bazli analitik: donusum orani, kaynak analizi"),
            ("\U0001F4B0", "Odeme plani olusturma, taksitlendirme, makbuz kesme"),
            ("\U0001F4E7", "Otomatik SMS/email bildirimleri (basvuru onay, hatirlatma)"),
        ],
        "PIPELINE ADIMLARI:\n\n"
        "\u25C6 1. On Basvuru\n"
        "   Web formu veya manuel giris\n\n"
        "\u25C6 2. Gorusme\n"
        "   Veli gorusmesi & rehberlik\n\n"
        "\u25C6 3. Sinav/Degerlendirme\n"
        "   Seviye tespit sinavi\n\n"
        "\u25C6 4. Kabul Karari\n"
        "   Komisyon degerlendirmesi\n\n"
        "\u25C6 5. Kesin Kayit\n"
        "   Evrak + Odeme + Onay"
    )


def slide_07_akademik():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4DA", "Akademik Takip", "Ogrenci basarisinin 360 derece takibi",
        [("30", "Sekme"), ("5", "Ana Grup"), ("9", "Veri Modeli"), ("7", "Rapor")],
        [
            ("\U0001F4D6", "Yoklama & devamsizlik: gunluk/haftalik/toplu giris, 7 alt sekme"),
            ("\U0001F4DD", "Not girisi: coklu not turu, donem bazli, otomatik ortalama"),
            ("\U0001F4D3", "Ders defteri: kazanim isleme takibi, ay bazli"),
            ("\U0001F4C5", "Akademik planlama: yillik/aylik/haftalik plan, MEB uyumu"),
            ("\U0001F4CA", "Karne & siralama: otomatik karne, sinif/okul siralamasi"),
            ("\U0001F4E6", "Odev takip: dosya/link/video/QR, online teslim"),
        ],
        "VERI MODELLERI:\n\n"
        "\u25C6 Student — Ogrenci kaydi\n"
        "\u25C6 Teacher — Ogretmen profili\n"
        "\u25C6 GradeRecord — Not kaydi\n"
        "\u25C6 AttendanceRecord — Devamsizlik\n"
        "\u25C6 ScheduleSlot — Ders programi\n"
        "\u25C6 AcademicPlan — Plan kaydi\n"
        "\u25C6 KazanimIsleme — Kazanim takip\n"
        "\u25C6 OlcmeTakvim — Sinav takvimi\n"
        "\u25C6 EtutKayit — Etut/destek kaydi"
    )


def slide_08_olcme():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4DD", "Olcme & Degerlendirme", "Turkiye'nin en gelismis sinav sistemi",
        [("42", "Sektor Gap"), ("7", "Soru Tipi"), ("20", "Blueprint"), ("6", "Bloom")],
        [
            ("\U0001F916", "AI soru uretimi: GPT-4o-mini, kazanim bazli, 5 adim wizard"),
            ("\U0001F4DA", "Soru bankasi: MCQ, dogru/yanlis, eslestirme, siralama, cloze, matematik"),
            ("\U0001F3AF", "Online sinav: tab guvenlik, heartbeat, negatif puanlama"),
            ("\U0001F4CA", "Psikometrik analiz: IRT (1PL, 2PL, 3PL), madde analizi"),
            ("\U0001F4C4", "OSYM tarzi PDF: 2 sutun, optik form, QR kod"),
            ("\U0001F9E9", "Adaptif test: CAT engine, Fisher bilgi, MLE"),
        ],
        "SINAV ALTYAPISI:\n\n"
        "\u25C6 ExamGenerator — Blueprint'ten sinav\n"
        "\u25C6 AutoGrader — Otomatik puanlama\n"
        "\u25C6 ProctoringEngine — Gozetim\n"
        "\u25C6 QTIExporter — IMS QTI 2.1\n"
        "\u25C6 LTIProvider — Moodle/Canvas\n"
        "\u25C6 IRTModel — Rasch modeli\n"
        "\u25C6 CATEngine — Adaptif test\n"
        "\u25C6 DifficultyCalibrator — p-value\n"
        "\u25C6 LearningAnalytics — Analitik"
    )


def slide_09_rehberlik():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F9E0", "Rehberlik & Psikolojik Danismanlik",
        "Ogrenci ruh sagligi ve gelisim takibi",
        [("8", "Sekme"), ("5", "Gorusme Tipi"), ("4", "Risk Seviye"), ("3", "BEP Alani")],
        [
            ("\U0001F4CB", "Vaka yonetimi: acma, takip, kapatma sureci"),
            ("\U0001F4AC", "Gorusme kaydi: ogrenci, veli, ogretmen gorusmeleri"),
            ("\U0001F60A", "Mood check-in: gunluk duygu durumu takibi"),
            ("\u26A0\uFE0F", "Risk analizi: erken uyari, zorbalik, dropout riski"),
            ("\U0001F4D1", "BEP takibi: bireysel egitim programi, hedef izleme"),
            ("\U0001F6A8", "Kriz mudahale: acil durum protokolleri"),
        ],
        "MOOD CHECK-IN:\n\n"
        "\U0001F60A Cok Mutlu\n"
        "\U0001F642 Mutlu\n"
        "\U0001F610 Normal\n"
        "\U0001F641 Uzgun\n"
        "\U0001F622 Cok Uzgun\n\n"
        "Gunluk takip ile erken mudahale.\n"
        "AI destekli risk tespiti.\n"
        "Anonim sinif bazli raporlama."
    )


def slide_10_ik():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F465", "IK & Personel Yonetimi",
        "Ise alimdan performansa tum surec",
        [("6", "Ana Sekme"), ("12", "Alt Sekme"), ("5", "Mulakat"), ("4", "Izin Turu")],
        [
            ("\U0001F4CB", "Ise alim: pozisyon acma, ilan, basvuru toplama"),
            ("\U0001F3AF", "Mulakat sureci: 5 asamali degerlendirme"),
            ("\U0001F4C5", "Izin yonetimi: yillik, mazeret, saglik, ucretsiz"),
            ("\U0001F4CA", "Performans: donemsel degerlendirme, hedef takibi"),
            ("\U0001F4B0", "Bordro entegrasyon: maas, prim, mesai hesaplama"),
            ("\U0001F4C4", "Ozluk dosyasi: dijital dokuman arsivi"),
        ],
        "ISE ALIM PIPELINE:\n\n"
        "\u25C6 Pozisyon Talebi\n"
        "\u25C6 Ilan Yayinlama\n"
        "\u25C6 Basvuru Toplama\n"
        "\u25C6 On Eleme\n"
        "\u25C6 Mulakat (5 asama)\n"
        "\u25C6 Referans Kontrolu\n"
        "\u25C6 Teklif & Ise Baslatma\n\n"
        "Tum surec dijital ve izlenebilir."
    )


def slide_11_butce():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4B0", "Butce & Odeme Takip",
        "Gelir/gider yonetimi ve veli odeme takibi",
        [("4", "Ana Modul"), ("12+", "Rapor"), ("3", "Odeme Tipi"), ("5", "Geciken Band")],
        [
            ("\U0001F4B8", "Gelir/gider takibi: kategori bazli, donem karsilastirma"),
            ("\U0001F4CB", "Taksit plani: ogrenci bazli, esnek taksitlendirme"),
            ("\U0001F9FE", "Makbuz kesme: otomatik numara, PDF cikti"),
            ("\U0001F514", "Geciken uyari: SMS/email hatirlatma, renk kodlu takip"),
            ("\U0001F4CA", "Butce planlama: yillik projeksiyon, sapma analizi"),
            ("\U0001F4B3", "Coklu odeme yontemi: nakit, havale, kredi karti"),
        ],
        "ODEME TAKIP OZELLIKLERI:\n\n"
        "\u25C6 Taksit plani olusturma\n"
        "\u25C6 Otomatik hatirlatma SMS\n"
        "\u25C6 Geciken odeme raporlari\n"
        "\u25C6 Makbuz/dekont PDF\n"
        "\u25C6 Donem bazli analiz\n"
        "\u25C6 Indirim & burs yonetimi\n"
        "\u25C6 Kardes indirimi otomatik\n"
        "\u25C6 Tahsilat dashboard"
    )


def slide_12_kurum():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F3EB", "Kurum Hizmetleri",
        "Gunluk operasyonlarin dijital yonetimi",
        [("7", "Hizmet Alani"), ("15+", "Sekme"), ("5", "Nobet Turu"), ("30", "Menu Gun")],
        [
            ("\U0001F6E1\uFE0F", "Nobet yonetimi: ogretmen/ogrenci nobeti, lokasyon bazli"),
            ("\u23F0", "Zaman cizelgesi: ders/teneffus/ogle/etut zamanlama"),
            ("\U0001F4C5", "Ders programi: sinif/ogretmen/salon bazli goruntuleme"),
            ("\U0001F35D", "Yemek menusu: haftalik menu, alerji takibi, diyet secenekleri"),
            ("\U0001F68C", "Servis takip: GPS canli konum, rota optimizasyonu"),
            ("\U0001F4E9", "Veli talepleri: dijital form, takip numarasi, durum bildirimi"),
        ],
        "YEMEK & ALERJI:\n\n"
        "\u25C6 Haftalik menu planlama\n"
        "\u25C6 Ogrenci bazli alerji kaydi\n"
        "\u25C6 Diyet secenekleri (vejeteryan, vegan)\n"
        "\u25C6 Kalori & besin degeri hesabi\n"
        "\u25C6 Veli bilgilendirme\n\n"
        "SERVIS GPS:\n\n"
        "\u25C6 Canli konum takibi\n"
        "\u25C6 Veli mobil erisim\n"
        "\u25C6 Tahmini varis suresi"
    )


def slide_13_uyari():
    slide = new_slide()
    _module_slide(slide,
        "\u26A0\uFE0F", "Erken Uyari Sistemi",
        "Yapay zeka destekli 32 boyutlu risk analizi",
        [("32", "Risk Boyutu"), ("4", "Risk Seviye"), ("3", "AI Model"), ("7/24", "Izleme")],
        [
            ("\U0001F916", "AI tahmin motoru: dropout riski, basari dususu, davranis degisimi"),
            ("\U0001F534", "4 seviye risk: Dusuk / Orta / Yuksek / Kritik"),
            ("\U0001F4CA", "32 boyutlu analiz: akademik, sosyal, devamsizlik, psikolojik..."),
            ("\U0001F514", "Otomatik bildirim: rehber, yonetici, veli bilgilendirme"),
            ("\U0001F9E9", "Zorbalik tespiti: davranis kaliplari analizi"),
            ("\U0001F4CB", "Mudahale plani: otomatik oneriler, takip sureci"),
        ],
        "RISK BOYUTLARI (Ornek):\n\n"
        "\u25C6 Akademik basari trendi\n"
        "\u25C6 Devamsizlik orani\n"
        "\u25C6 Not degisimi hizi\n"
        "\u25C6 Sosyal etkilesim skoru\n"
        "\u25C6 Mood check-in trendi\n"
        "\u25C6 Odev teslim orani\n"
        "\u25C6 Disiplin kayitlari\n"
        "\u25C6 Aile durumu degisimi\n"
        "\u25C6 Akran iliskileri\n"
        "\u25C6 ... ve 23 boyut daha"
    )


def slide_14_veli():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F468\u200D\U0001F469\u200D\U0001F467", "Veli Paneli",
        "Veliler icin ozel tasarlanmis mobil deneyim",
        [("5", "Cocuk Sekmesi"), ("8", "Kapsul"), ("3", "Randevu Tipi"), ("1-Tik", "Odeme")],
        [
            ("\U0001F476", "Cocuk detay: 5 sekme (genel, notlar, devamsizlik, odevler, gorusmeler)"),
            ("\U0001F4CA", "Kapsul gorunum: ozet kartlar ile hizli erisim"),
            ("\U0001F4C5", "Randevu: online randevu alma, takvim goruntuleme"),
            ("\U0001F68C", "Servis: GPS canli takip, binis/inis bildirimi"),
            ("\U0001F4B0", "Odeme: taksit durumu, gecmis odemeler, online odeme"),
            ("\U0001F4AC", "Mesajlasma: ogretmen ile direkt iletisim"),
        ],
        "VELI MOBIL DENEYIMI:\n\n"
        "\u25C6 Push bildirimler\n"
        "\u25C6 Cocuk secimi (coklu cocuk)\n"
        "\u25C6 Haftalik ozet rapor\n"
        "\u25C6 Sinav sonuc bildirimi\n"
        "\u25C6 Devamsizlik uyarisi\n"
        "\u25C6 Menu goruntuleme\n"
        "\u25C6 Servis GPS canli\n"
        "\u25C6 Odeme hatirlatma\n"
        "\u25C6 Duyuru & etkinlikler"
    )


def slide_15_ogrenci():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F393", "Ogrenci Paneli",
        "Ogrenciler icin motive edici dijital deneyim",
        [("24", "Sayfa"), ("5", "Dashboard"), ("3", "AI Ozellik"), ("10+", "Oyun")],
        [
            ("\U0001F4CA", "Kisisel dashboard: not ortalamasi, siralama, hedefler"),
            ("\U0001F4D3", "Dijital defterim: ders notlari, odevler, kaynaklar"),
            ("\U0001F916", "AI asistan: soru sorma, aciklama, ders plani onerisi"),
            ("\U0001F3AE", "Egitici oyunlar: matematik, kelime, hafiza, bulmaca"),
            ("\U0001F3C6", "Basari rozetleri: gamification ile motivasyon"),
            ("\U0001F4DA", "Online sinav: zamana karsi yaris, aninda sonuc"),
        ],
        "GAMIFICATION:\n\n"
        "\u25C6 XP puan sistemi\n"
        "\u25C6 Seviye atlama\n"
        "\u25C6 Basari rozetleri\n"
        "\u25C6 Haftalik liderlik tablosu\n"
        "\u25C6 Gorev tamamlama oduller\n\n"
        "AI ASISTAN:\n\n"
        "\u25C6 7/24 soru yanit\n"
        "\u25C6 Ders anlatimi\n"
        "\u25C6 Calisma plani onerisi"
    )


def slide_16_dijital():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F310", "Dijital Ogrenme",
        "5 dil, CEFR, STEAM ve AI destekli ogrenme",
        [("5", "Dil"), ("6", "CEFR Seviye"), ("4", "STEAM Alani"), ("1", "AI Treni")],
        [
            ("\U0001F1EC\U0001F1E7", "Ingilizce, Almanca, Fransizca, Ispanyolca, Arapca"),
            ("\U0001F4CA", "CEFR seviye belirleme: A1-C2 otomatik yerlestirme"),
            ("\U0001F916", "AI Treni: interaktif, gomulu, backend gerektirmez"),
            ("\U0001F52C", "STEAM projeleri: bilim, teknoloji, muhendislik, sanat, matematik"),
            ("\U0001F3AE", "Interaktif alistirmalar: dinleme, konusma, yazma, okuma"),
            ("\U0001F4CB", "Ilerleme takibi: beceri bazli, haftalik rapor"),
        ],
        "DESTEKLENEN DILLER:\n\n"
        "\U0001F1EC\U0001F1E7 Ingilizce (Cambridge)\n"
        "\U0001F1E9\U0001F1EA Almanca (Goethe)\n"
        "\U0001F1EB\U0001F1F7 Fransizca (DELF)\n"
        "\U0001F1EA\U0001F1F8 Ispanyolca (DELE)\n"
        "\U0001F1F8\U0001F1E6 Arapca\n\n"
        "AI TRENI:\n"
        "Tamamen gomulu calisir.\n"
        "Backend gerektirmez.\n"
        "Interaktif soru-cevap."
    )


def slide_17_kutuphane():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4DA", "Kutuphane & Etkinlik",
        "Barkodlu kutuphane + kulup & etkinlik yonetimi",
        [("3", "Kutuphane Mod"), ("5", "Kulup Tipi"), ("12", "Etkinlik/Yil"), ("1-Tik", "Barkod")],
        [
            ("\U0001F4D6", "Kitap katalogu: barkod ile hizli kayit, ISBN arama"),
            ("\U0001F504", "Odunc/iade: barkod okutma, otomatik sure takibi"),
            ("\U0001F514", "Gecikme uyari: SMS/email, ceza hesaplama"),
            ("\U0001F3AD", "Kulup yonetimi: kayit, yoklama, etkinlik planlama"),
            ("\U0001F3C6", "Etkinlik takvimi: okul geneli, sinif bazli, disaridan"),
            ("\U0001F4CA", "Okuma istatistikleri: ogrenci bazli, sinif karsilastirma"),
        ],
        "BARKOD SISTEMI:\n\n"
        "\u25C6 ISBN barkod okuma\n"
        "\u25C6 Ogrenci kimlik barkod\n"
        "\u25C6 1 saniyede odunc/iade\n"
        "\u25C6 Stok sayim kolayligi\n\n"
        "KULUP YONETIMI:\n\n"
        "\u25C6 Spor kulupler\n"
        "\u25C6 Sanat & muzik\n"
        "\u25C6 Bilim & teknoloji\n"
        "\u25C6 Sosyal sorumluluk\n"
        "\u25C6 Liderlik programi"
    )


def slide_18_saglik():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F3E5", "Saglik & Guvenlik",
        "Revir, sivil savunma, is sagligi guvenligi",
        [("3", "Ana Alan"), ("5", "Revir Kayit"), ("4", "Tatbikat/Yil"), ("7/24", "Kriz Hatti")],
        [
            ("\U0001F48A", "Revir kayitlari: saglik gecmisi, ilac takibi, alerji"),
            ("\U0001F691", "Acil durum: kriz protokolu, ambulans cagrisi, veli bilgilendirme"),
            ("\U0001F6E1\uFE0F", "Sivil savunma: deprem/yangin tatbikati, tahliye plani"),
            ("\U0001F9F0", "ISG: risk degerlendirme, is kazasi kaydi, denetim"),
            ("\U0001F489", "Asi takibi: ogrenci asi kayitlari, hatirlatma"),
            ("\U0001F4CB", "Saglik raporu: donemsel saglik istatistikleri"),
        ],
        "REVIR KAYIT SISTEMI:\n\n"
        "\u25C6 Ogrenci saglik profili\n"
        "\u25C6 Kronik hastalik takibi\n"
        "\u25C6 Ilac bilgisi (alerji!)\n"
        "\u25C6 Revir ziyaret kaydi\n"
        "\u25C6 Kaza/yaralanma formu\n\n"
        "SIVIL SAVUNMA:\n\n"
        "\u25C6 Deprem tatbikati plani\n"
        "\u25C6 Yangin tahliye rotasi\n"
        "\u25C6 Toplanma alani bilgisi\n"
        "\u25C6 Veli bilgilendirme SMS"
    )


def slide_19_analitik():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4CA", "Analitik Dashboard",
        "Plotly ile gorsellestirme, veri odakli karar destek",
        [("5", "Dashboard Tab"), ("15+", "Grafik Tipi"), ("3", "Karsilastirma"), ("PDF", "Export")],
        [
            ("\U0001F4CA", "Sinif bazli basari analizi: ortalama, dagılım, trend"),
            ("\U0001F4C8", "Donem karsilastirma: 1. ve 2. donem, yillik trend"),
            ("\U0001F3AF", "Ders bazli analiz: en basarili/basarisiz dersler"),
            ("\U0001F465", "Ogretmen performansi: sinif ortalamalari karsilastirma"),
            ("\U0001F4CB", "Devamsizlik analizi: gun/saat bazli, trend grafikleri"),
            ("\U0001F4E5", "PDF/Excel export: tum raporlar indirilebilir"),
        ],
        "PLOTLY GORSELLESTIRME:\n\n"
        "\u25C6 Bar chart — sinif karsilastirma\n"
        "\u25C6 Line chart — trend analizi\n"
        "\u25C6 Pie chart — dagılım\n"
        "\u25C6 Heatmap — ders/sinif matrix\n"
        "\u25C6 Box plot — not dagılımı\n"
        "\u25C6 Scatter — korelasyon\n"
        "\u25C6 Radar — coklu boyut\n\n"
        "Tum grafikler interaktif,\n"
        "zoom + hover detay."
    )


def slide_20_sosyal():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F4E2", "Sosyal Medya & Iletisim",
        "Icerik takvimi, mesajlasma, gorusme yonetimi",
        [("4", "Kanal"), ("30", "Gun/Ay Plan"), ("3", "Mesaj Turu"), ("2", "SMS Saglayici")],
        [
            ("\U0001F4C5", "Icerik takvimi: aylik plan, onay sureci, otomatik paylasim"),
            ("\U0001F4AC", "Mesajlasma: veli-ogretmen, yonetici duyuru, toplu mesaj"),
            ("\U0001F4E7", "Email: SMTP entegrasyon, sablon destekli, toplu gonderim"),
            ("\U0001F4F1", "SMS: NetGSM/Twilio, otomatik hatirlatma, maliyet takibi"),
            ("\U0001F4CB", "Gorusme kaydi: tarih, katilimci, konu, sonuc"),
            ("\U0001F4CA", "Iletisim analitigi: acilma orani, yanit suresi"),
        ],
        "ILETISIM KANALLARI:\n\n"
        "\u25C6 SMS (NetGSM / Twilio)\n"
        "\u25C6 Email (SMTP)\n"
        "\u25C6 Push Bildirim (Mobil)\n"
        "\u25C6 Uygulama Ici Mesaj\n\n"
        "OTOMATIK BILDIRIMLER:\n\n"
        "\u25C6 Sinav sonucu\n"
        "\u25C6 Devamsizlik uyarisi\n"
        "\u25C6 Odeme hatirlatma\n"
        "\u25C6 Etkinlik duyurusu\n"
        "\u25C6 Randevu onay/iptal"
    )


def slide_21_sertifika():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F3C6", "Sertifika Uretici",
        "7 sablon, toplu PDF uretimi, ZIP export",
        [("7", "Sablon"), ("PDF", "Cikti"), ("ZIP", "Toplu"), ("QR", "Dogrulama")],
        [
            ("\U0001F4DC", "7 profesyonel sablon: katilim, basari, tesekkur, spor, sanat..."),
            ("\U0001F3A8", "Kurumsal tasarim: logo, renkler, imza alani"),
            ("\U0001F4E5", "Toplu uretim: sinif bazli, etkinlik bazli, ZIP indirme"),
            ("\U0001F4F1", "QR dogrulama: her sertifikada benzersiz QR kod"),
            ("\U0001F58A\uFE0F", "Dijital imza: yonetici + ogretmen imza alani"),
            ("\U0001F4CA", "Arsiv: tum uretilen sertifikalar kayit altinda"),
        ],
        "SABLON TURLERI:\n\n"
        "\u25C6 Katilim Belgesi\n"
        "\u25C6 Basari Belgesi\n"
        "\u25C6 Tesekkur Belgesi\n"
        "\u25C6 Spor Basari\n"
        "\u25C6 Sanat & Muzik\n"
        "\u25C6 Bilim Olimpiyadi\n"
        "\u25C6 Ozel Tasarim\n\n"
        "Her sertifika PDF formatinda,\n"
        "QR kod ile dogrulanabilir."
    )


def slide_22_destek():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F527", "Destek Hizmetleri",
        "Bakim, talep, denetim ve tesis yonetimi",
        [("4", "Ana Alan"), ("8", "Talep Turu"), ("3", "Oncelik"), ("5", "Denetim")],
        [
            ("\U0001F4CB", "Talep yonetimi: ariza bildirimi, oncelik, atama, takip"),
            ("\U0001F527", "Bakim planlama: periyodik bakim takvimi, hatirlatma"),
            ("\U0001F50D", "Denetim: bina/sinif denetimi, kontrol listesi, puanlama"),
            ("\U0001F3EB", "Tesis yonetimi: salon/sinif envanteri, kapasite planlama"),
            ("\U0001F4CA", "Raporlama: ariza istatistikleri, cozum suresi analizi"),
            ("\U0001F514", "Bildirim: talep durumu degisikliginde otomatik uyari"),
        ],
        "TALEP SURECI:\n\n"
        "\u25C6 1. Bildirim (herkes)\n"
        "\u25C6 2. Onceliklendirme\n"
        "\u25C6 3. Atama (teknik ekip)\n"
        "\u25C6 4. Is emri olusturma\n"
        "\u25C6 5. Cozum & onay\n"
        "\u25C6 6. Kapatma & rapor\n\n"
        "Ortalama cozum suresi: <24 saat\n"
        "Memnuniyet anketi: her talep sonrasi"
    )


def slide_23_toplanti():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F91D", "Toplanti & Randevu",
        "Ogretmenler kurulu, ziyaretci yonetimi",
        [("3", "Toplanti Tipi"), ("6", "Randevu Tab"), ("5", "Ziyaretci Adim"), ("PDF", "Tutanak")],
        [
            ("\U0001F4CB", "Ogretmenler kurulu: gundem, katilim, tutanak, karar takibi"),
            ("\U0001F4C5", "Randevu sistemi: online alma, onaylama, hatirlatma"),
            ("\U0001F6B6", "Ziyaretci giris/cikis: kimlik kaydi, ziyaret nedeni"),
            ("\U0001F4DC", "Toplanti tutanagi: otomatik PDF, katilimci imzasi"),
            ("\U0001F514", "Hatirlatma: SMS/email ile toplanti oncesi uyari"),
            ("\U0001F4CA", "Istatistik: toplanti suresi, katilim orani analizi"),
        ],
        "ZIYARETCI YONETIMI:\n\n"
        "\u25C6 1. Kimlik kaydi\n"
        "\u25C6 2. Ziyaret nedeni\n"
        "\u25C6 3. Gorusulecek kisi\n"
        "\u25C6 4. Giris saati kaydi\n"
        "\u25C6 5. Cikis saati kaydi\n\n"
        "Gorusulecek unvanlar IK\n"
        "modulunden otomatik gelir.\n"
        "Tum ziyaretler arsivlenir."
    )


def slide_24_mezun():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F393", "Mezunlar & Kariyer",
        "Mezun takibi ve universite yerlestirme",
        [("3", "Ana Alan"), ("5", "Takip Kriteri"), ("100+", "Universite"), ("PDF", "Rapor")],
        [
            ("\U0001F393", "Mezun veritabani: iletisim, universite, kariyer bilgileri"),
            ("\U0001F3EB", "Universite yerlestirme: YKS sonuc takibi, istatistikler"),
            ("\U0001F4BC", "Kariyer rehberligi: ilgi alani testi, meslek onerileri"),
            ("\U0001F4CA", "Basari haritasi: mezunlarin universite dagılımı"),
            ("\U0001F4AC", "Mezun iletisim: etkinlik davet, referans programa"),
            ("\U0001F3C6", "Basari hikayeleri: vitrin sayfasi, motivasyon"),
        ],
        "UNIVERSITE ISTATISTIKLERI:\n\n"
        "\u25C6 YKS puan dagılımı\n"
        "\u25C6 Bolum bazli yerlestirme\n"
        "\u25C6 Yillara gore trend\n"
        "\u25C6 Turkiye siralamasi\n"
        "\u25C6 Burslu ogrenci orani\n\n"
        "KARIYER REHBERLIGI:\n\n"
        "\u25C6 Holland ilgi envanteti\n"
        "\u25C6 AI bazli oneri motoru\n"
        "\u25C6 Mentor eslestirme"
    )


def slide_25_ai():
    slide = new_slide()
    _module_slide(slide,
        "\U0001F916", "SmartI — AI Asistan",
        "GPT-4o-mini destekli yapay zeka copilot",
        [("GPT-4o", "AI Model"), ("7/24", "Erisim"), ("5", "Mod"), ("3", "Dil")],
        [
            ("\U0001F4AC", "Sohbet modu: ogrenci/ogretmen/veli icin AI chat"),
            ("\U0001F399\uFE0F", "Sesli mod: konusarak soru sorma, AI yanit"),
            ("\U0001F4D6", "Ders plani copilot: AI ile haftalik plan olusturma"),
            ("\U0001F4DD", "Soru uretimi: kazanim bazli, Bloom taksonomisi, 7 soru tipi"),
            ("\U0001F4CA", "Analiz: ogrenci performans ozeti, oneri motoru"),
            ("\U0001F30D", "Coklu dil: Turkce, Ingilizce, Almanca destek"),
        ],
        "AI KULLANIM ALANLARI:\n\n"
        "\u25C6 Soru bankasi uretimi\n"
        "\u25C6 PDF'den soru cikarma\n"
        "\u25C6 Ders plani olusturma\n"
        "\u25C6 Ogrenci analiz raporu\n"
        "\u25C6 Telafi gorevi olusturma\n"
        "\u25C6 Lead skorlama\n"
        "\u25C6 Mood analizi\n"
        "\u25C6 Risk tahmini\n\n"
        "Tum AI islemleri audit log ile\n"
        "kayit altindadir."
    )


# ====================================================================
# MOBILE SLIDES (26–28)
# ====================================================================

def slide_26_mobil():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "MOBIL UYGULAMA")
    _add_rich_text_box(
        slide, MARGIN_LEFT, Inches(0.55), CONTENT_WIDTH, Inches(0.6),
        [
            ("\U0001F4F1  ", Pt(30), GOLD, False, "Segoe UI Emoji"),
            ("SmartCampus Mobil", Pt(30), WHITE, True, "Segoe UI Black"),
        ]
    )
    _add_line(slide, MARGIN_LEFT, Inches(1.2), Inches(5), Inches(1.2),
              color=GOLD, width=Pt(3))
    _add_text_box(slide, MARGIN_LEFT, Inches(1.35), Inches(7), Inches(0.35),
                  "Flutter Android — 99 Dart dosyasi, 5 rol, ultra premium tasarim",
                  Pt(14), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)

    # Stats
    stats = [("99", "Dart Dosyasi"), ("5", "Rol"), ("50+", "Ekran"), ("Dark/Light", "Tema")]
    sx = MARGIN_LEFT
    for num, label in stats:
        _add_stat_card(slide, sx, Inches(1.9), Inches(2.5), Inches(1.0),
                       num, label, border_color=DARK_GOLD)
        sx += Inches(2.7)

    # Feature columns
    left_feats = [
        ("\U0001F451", "Yonetici: Dashboard, ogrenci listesi, bildirim"),
        ("\U0001F4DA", "Ogretmen: Siniflarim, yoklama, not, ders defteri"),
        ("\U0001F9E0", "Rehber: Vaka, gorusme, mood, risk"),
    ]
    right_feats = [
        ("\U0001F468\u200D\U0001F469\u200D\U0001F467", "Veli: Cocuk takip, randevu, servis, odeme"),
        ("\U0001F393", "Ogrenci: Dashboard, AI, sinav, oyunlar"),
        ("\U0001F3A8", "Ultra Premium: Gold-navy tasarim, glass effect, shimmer"),
    ]

    _add_icon_bullet_list(slide, MARGIN_LEFT, Inches(3.2),
                          Inches(5.5), Inches(3.0), left_feats,
                          font_size=Pt(12), color=OFF_WHITE, spacing=1.5)
    _add_icon_bullet_list(slide, Inches(7.0), Inches(3.2),
                          Inches(5.5), Inches(3.0), right_feats,
                          font_size=Pt(12), color=OFF_WHITE, spacing=1.5)


def slide_27_mobil_ozellik():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "MOBIL OZELLIKLER")
    _slide_title(slide, "Mobil Uygulama Detay", title_y=Inches(0.5))

    features = [
        ("\U0001F4F6", "Offline Cache", "Internet olmadan da calisir.\nVeri senkronizasyonu otomatik.", CYAN),
        ("\U0001F514", "Push Bildirim", "Aninda bildirim.\nSinav, not, duyuru, randevu.", GOLD),
        ("\U0001F3A8", "Dark/Light Tema", "Kullanici tercihine gore.\nGoz yorgunlugu azaltici.", PURPLE),
        ("\U0001F512", "Sifre Degistirme", "Guvenli giris.\nbcrypt sifreleme.", GREEN),
        ("\U0001F4F1", "Responsive UI", "Her ekran boyutunda\nmukemmel gorunum.", BLUE),
        ("\U0001F504", "Auto Update", "Yeni versiyon bildirimi.\nTek tikla guncelleme.", ORANGE),
    ]

    cw = Inches(3.6)
    ch = Inches(2.0)
    gap = Inches(0.3)
    sx = Inches(0.7)
    sy = Inches(1.8)

    for i, (emoji, title, desc, accent) in enumerate(features):
        row = i // 3
        col = i % 3
        x = sx + col * (cw + gap)
        y = sy + row * (ch + gap)
        _feature_card(slide, x, y, cw, ch, emoji, title, desc,
                      border_color=accent)


def slide_28_apk():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "CI/CD & DAGITIM")
    _slide_title(slide, "APK Build & Dagitim", title_y=Inches(0.5))

    # Pipeline
    steps = [
        ("1\uFE0F\u20E3", "Git Push", "Kod degisikligi\nGitHub'a push"),
        ("2\uFE0F\u20E3", "GitHub Actions", "Otomatik CI/CD\nTest & build"),
        ("3\uFE0F\u20E3", "Flutter Build", "APK uretimi\nRelease mode"),
        ("4\uFE0F\u20E3", "Play Store", "Dagitim hazir\nAAB format"),
    ]

    sw = Inches(2.5)
    sh = Inches(2.0)
    gap = Inches(0.5)
    sx = Inches(0.9)
    sy = Inches(1.8)

    for i, (num, title, desc) in enumerate(steps):
        x = sx + i * (sw + gap)
        _add_rounded_rect(slide, x, sy, sw, sh,
                          fill_color=CARD_BG, border_color=GOLD,
                          border_width=Pt(1))
        _add_text_box(slide, x, sy + Inches(0.15), sw, Inches(0.4),
                      num, Pt(28), GOLD, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x, sy + Inches(0.6), sw, Inches(0.35),
                      title, Pt(14), WHITE, PP_ALIGN.CENTER, bold=True)
        _add_text_box(slide, x + Inches(0.2), sy + Inches(1.0),
                      sw - Inches(0.4), Inches(0.8),
                      desc, Pt(11), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                      line_spacing=1.4)

        # Arrow between steps
        if i < 3:
            arrow_x = x + sw + Inches(0.05)
            _add_text_box(slide, arrow_x, sy + Inches(0.7),
                          Inches(0.4), Inches(0.4),
                          "\u25B6", Pt(20), GOLD, PP_ALIGN.CENTER, bold=False)

    # Bottom info
    _add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.2),
                      fill_color=CARD_BG_ALT, border_color=DARK_GOLD)
    info_items = [
        ("\u2705", "GitHub Actions workflow: otomatik test + build her push'ta"),
        ("\u2705", "Flutter 3.x + Dart 3.x: en guncel framework"),
        ("\u2705", "Release APK: ProGuard optimizasyonu, ~25MB boyut"),
        ("\u2705", "AAB format: Play Store icin optimize edilmis"),
        ("\u2705", "Versiyon yonetimi: semantic versioning (major.minor.patch)"),
        ("\u2705", "Kod imzalama: keystore ile guvenli dagitim"),
    ]
    _add_icon_bullet_list(slide, Inches(1.1), Inches(4.5),
                          Inches(11.0), Inches(1.8), info_items,
                          font_size=Pt(11), color=OFF_WHITE, spacing=1.3)


# ====================================================================
# TECHNICAL SLIDES (29–31)
# ====================================================================

def slide_29_teknik():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "TEKNIK ALTYAPI")
    _slide_title(slide, "Mimari & Teknoloji Stack", title_y=Inches(0.5))

    # Tech stack cards
    techs = [
        ("\U0001F4BB", "FRONTEND", "Streamlit (Python)", [
            "30+ modul, 390+ sekme",
            "Responsive tasarim",
            "Custom CSS temaları",
            "Session state yonetimi",
        ], BLUE),
        ("\U0001F4F1", "MOBILE", "Flutter (Dart)", [
            "99 sayfa, 5 rol",
            "Material Design 3",
            "Provider state mgmt",
            "Offline-first mimari",
        ], GREEN),
        ("\u2699\uFE0F", "BACKEND", "FastAPI (Python)", [
            "143 REST API route",
            "Async/await pattern",
            "Pydantic validation",
            "OpenAPI (Swagger) docs",
        ], PURPLE),
        ("\U0001F4BE", "DATA", "JSON + File Store", [
            "687+ veri dosyasi",
            "Multi-tenant izolasyon",
            "Otomatik yedekleme",
            "SQLite'a gecise hazir",
        ], ORANGE),
    ]

    tw = Inches(2.7)
    th = Inches(4.0)
    gap = Inches(0.3)
    sx = Inches(0.7)
    sy = Inches(1.7)

    for i, (emoji, layer, tech, items, accent) in enumerate(techs):
        x = sx + i * (tw + gap)
        _add_rounded_rect(slide, x, sy, tw, th,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1.5))
        _add_text_box(slide, x, sy + Inches(0.15), tw, Inches(0.4),
                      emoji, Pt(28), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x, sy + Inches(0.55), tw, Inches(0.3),
                      layer, Pt(13), WHITE, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Tech badge
        bw = Inches(2.0)
        _add_rounded_rect(slide, x + (tw - bw) / 2, sy + Inches(0.9),
                          bw, Inches(0.28), fill_color=accent)
        _add_text_box(slide, x + (tw - bw) / 2, sy + Inches(0.9),
                      bw, Inches(0.28),
                      tech, Pt(9), WHITE, PP_ALIGN.CENTER, bold=True)
        # Feature list
        for j, item in enumerate(items):
            _add_text_box(slide, x + Inches(0.15), sy + Inches(1.4) + j * Inches(0.5),
                          tw - Inches(0.3), Inches(0.4),
                          "\u25C6  " + item, Pt(10), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)

    # Bottom bar
    _add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.0), Inches(0.3),
                  "Python 3.11+  |  Dart 3.x  |  FastAPI 0.100+  |  Streamlit 1.30+  |  OpenAI API",
                  Pt(10), MED_GRAY, PP_ALIGN.CENTER, bold=False)


def slide_30_guvenlik():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "GUVENLIK")
    _slide_title(slide, "Enterprise Guvenlik", title_y=Inches(0.5))

    security = [
        ("\U0001F512", "JWT Authentication",
         "JSON Web Token ile guvenli oturum yonetimi. Access + refresh token pattern.",
         BLUE),
        ("\U0001F510", "bcrypt Sifreleme",
         "Endstri standardi sifre hashleme. Brute-force koruması. Salt + hash.",
         GREEN),
        ("\u26A1", "Rate Limiting",
         "API isteklerinde hiz sinirlandirma. DDoS korumasi. IP bazli limit.",
         ORANGE),
        ("\U0001F6E1\uFE0F", "CORS Policy",
         "Cross-origin guvenlik. Sadece yetkili domain'ler erisebilir.",
         PURPLE),
        ("\U0001F510", "Env Var Sifreler",
         "Hassas bilgiler .env dosyasinda. Kod icerisinde sifre yok.",
         RED),
        ("\U0001F465", "Rol Bazli Erisim",
         "5 rol, granular yetki. Her kullanici sadece yetkili ekranlari gorur.",
         CYAN),
    ]

    cw = Inches(3.6)
    ch = Inches(1.8)
    gap = Inches(0.3)
    sx = Inches(0.7)
    sy = Inches(1.8)

    for i, (emoji, title, desc, accent) in enumerate(security):
        row = i // 3
        col = i % 3
        x = sx + col * (cw + gap)
        y = sy + row * (ch + gap)
        _feature_card(slide, x, y, cw, ch, emoji, title, desc,
                      border_color=accent)


def slide_31_entegrasyon():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "ENTEGRASYONLAR")
    _slide_title(slide, "Dis Sistem Entegrasyonlari", title_y=Inches(0.5))

    integrations = [
        ("\U0001F4F1", "SMS", "NetGSM / Twilio",
         "Otomatik bildirim, toplu SMS, maliyet takibi, sablon yonetimi",
         BLUE),
        ("\U0001F4E7", "Email", "SMTP Protokolu",
         "Otomatik email, HTML sablon, toplu gonderim, bounce takibi",
         GREEN),
        ("\U0001F916", "AI", "OpenAI GPT-4o-mini",
         "Soru uretimi, ders plani, ogrenci analiz, chat asistan",
         PURPLE),
        ("\U0001F4B3", "Odeme", "Banka Entegrasyon",
         "Online odeme, taksit takibi, otomatik eslestirme, makbuz",
         GOLD),
    ]

    iw = Inches(5.5)
    ih = Inches(1.1)
    gap = Inches(0.2)
    sx = Inches(0.8)
    sy = Inches(1.8)

    for i, (emoji, name, tech, desc, accent) in enumerate(integrations):
        y = sy + i * (ih + gap)
        _add_rounded_rect(slide, sx, y, iw, ih,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1))
        _add_text_box(slide, sx + Inches(0.15), y + Inches(0.1),
                      Inches(0.4), Inches(0.4),
                      emoji, Pt(24), WHITE, PP_ALIGN.LEFT, bold=False,
                      font_name="Segoe UI Emoji")
        _add_rich_text_box(
            slide, sx + Inches(0.6), y + Inches(0.1), Inches(2.5), Inches(0.35),
            [
                (name, Pt(14), WHITE, True),
                ("  |  ", Pt(10), MED_GRAY, False),
                (tech, Pt(11), accent, True),
            ]
        )
        _add_text_box(slide, sx + Inches(0.6), y + Inches(0.5),
                      iw - Inches(0.8), Inches(0.5),
                      desc, Pt(10), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)

    # Right side — integration diagram concept
    _add_rounded_rect(slide, Inches(7.0), Inches(1.8),
                      Inches(5.5), Inches(4.5),
                      fill_color=CARD_BG_ALT, border_color=DARK_GOLD)
    _add_text_box(slide, Inches(7.2), Inches(1.95),
                  Inches(5.1), Inches(0.35),
                  "ENTEGRASYON AKISI", Pt(13), GOLD, PP_ALIGN.CENTER, bold=True)

    flow_items = [
        "\u25C6  SmartCampus Web/Mobil",
        "        \u2193",
        "\u25C6  FastAPI Backend (143 route)",
        "        \u2193",
        "\u25C6  Service Layer",
        "   \u251C\u2500 SMS Gateway (NetGSM/Twilio)",
        "   \u251C\u2500 Email Server (SMTP)",
        "   \u251C\u2500 OpenAI API (GPT-4o-mini)",
        "   \u2514\u2500 Payment Gateway",
        "        \u2193",
        "\u25C6  Data Store (JSON / SQLite)",
    ]
    _add_multiline_text(
        slide, Inches(7.3), Inches(2.4), Inches(4.9), Inches(3.5),
        [(line, Pt(10), OFF_WHITE if "\u25C6" in line else LIGHT_GRAY,
          "\u25C6" in line) for line in flow_items],
        alignment=PP_ALIGN.LEFT, line_spacing=1.3
    )


# ====================================================================
# COMPETITIVE SLIDES (32–34)
# ====================================================================

def slide_32_rakip():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "REKABET ANALIZI")
    _slide_title(slide, "Rakip Karsilastirma", title_y=Inches(0.5))

    # Comparison table
    headers = [
        ("OZELLIK", WHITE),
        ("SmartCampus AI", GOLD),
        ("Vedubox", LIGHT_GRAY),
        ("Edval", LIGHT_GRAY),
        ("Turkcell Kampus", LIGHT_GRAY),
    ]

    rows = [
        [("Web Modul", OFF_WHITE), ("30+  \u2705", GOLD), ("10-15", LIGHT_GRAY), ("5-8", LIGHT_GRAY), ("8-12", LIGHT_GRAY)],
        [("Mobil Uygulama", OFF_WHITE), ("5 Rol  \u2705", GOLD), ("1 Rol", LIGHT_GRAY), ("\u274C", RED), ("2 Rol", LIGHT_GRAY)],
        [("AI Soru Uretimi", OFF_WHITE), ("GPT-4o  \u2705", GOLD), ("\u274C", RED), ("\u274C", RED), ("\u274C", RED)],
        [("Online Sinav", OFF_WHITE), ("Tam  \u2705", GOLD), ("Temel", LIGHT_GRAY), ("\u274C", RED), ("Temel", LIGHT_GRAY)],
        [("Erken Uyari", OFF_WHITE), ("32 Boyut  \u2705", GOLD), ("\u274C", RED), ("\u274C", RED), ("\u274C", RED)],
        [("Servis GPS", OFF_WHITE), ("Canli  \u2705", GOLD), ("\u274C", RED), ("\u274C", RED), ("Var", LIGHT_GRAY)],
        [("Odeme Takip", OFF_WHITE), ("Tam  \u2705", GOLD), ("Temel", LIGHT_GRAY), ("\u274C", RED), ("Temel", LIGHT_GRAY)],
        [("Mood Check-in", OFF_WHITE), ("\u2705", GOLD), ("\u274C", RED), ("\u274C", RED), ("\u274C", RED)],
        [("Multi-tenant", OFF_WHITE), ("\u2705", GOLD), ("\u2705", GREEN), ("\u274C", RED), ("\u2705", GREEN)],
        [("Sertifika PDF", OFF_WHITE), ("7 Sablon  \u2705", GOLD), ("\u274C", RED), ("\u274C", RED), ("1 Sablon", LIGHT_GRAY)],
    ]

    table_left = Inches(0.5)
    table_top = Inches(1.7)
    col_w = Inches(12.333) / 5
    row_h = Inches(0.42)

    # Header row
    for ci, (text, color) in enumerate(headers):
        x = table_left + ci * col_w
        bg = DARK_GOLD if ci == 1 else CARD_BG_ALT
        _add_shape(slide, MSO_SHAPE.RECTANGLE, x, table_top, col_w, row_h,
                   fill_color=bg, border_color=DARK_GRAY, border_width=Pt(0.5))
        _add_text_box(slide, x + Inches(0.05), table_top, col_w - Inches(0.1), row_h,
                      text, Pt(10), color, PP_ALIGN.CENTER, bold=True)

    # Data rows
    for ri, row in enumerate(rows):
        y = table_top + (ri + 1) * row_h
        for ci, (text, color) in enumerate(row):
            x = table_left + ci * col_w
            bg = CARD_BG if ri % 2 == 0 else CARD_BG_ALT
            if ci == 1:
                bg = RGBColor(0x1A, 0x25, 0x37)  # slightly highlighted column
            _add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, col_w, row_h,
                       fill_color=bg, border_color=DARK_GRAY, border_width=Pt(0.5))
            _add_text_box(slide, x + Inches(0.05), y, col_w - Inches(0.1), row_h,
                          text, Pt(9), color, PP_ALIGN.CENTER, bold=(ci == 1))

    # Bottom note
    _add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.0), Inches(0.3),
                  "\u2705 = Tam ozellik  |  Temel = Sinirli  |  \u274C = Yok   —   SmartCampus AI her kategoride lider",
                  Pt(10), GOLD, PP_ALIGN.CENTER, bold=True)


def slide_33_fark():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "REKABET AVANTAJI")
    _slide_title(slide, "Fark Yaratan 10 Ozellik", title_y=Inches(0.5))

    features = [
        ("\U0001F916", "AI Soru Uretimi", "Kazanim bazli GPT-4o"),
        ("\U0001F60A", "Mood Check-in", "Gunluk duygu takibi"),
        ("\U0001F68C", "Servis GPS", "Canli konum takibi"),
        ("\U0001F4B0", "Odeme Takip", "Taksit + makbuz + geciken"),
        ("\u26A0\uFE0F", "32 Boyut Risk", "AI erken uyari sistemi"),
        ("\U0001F4DA", "Soru Bankasi", "7 tip, IRT, CAT, QTI"),
        ("\U0001F3C6", "Sertifika", "7 sablon, QR, toplu PDF"),
        ("\U0001F4D6", "Barkod Kutuphane", "1 saniye odunc/iade"),
        ("\U0001F4F1", "5 Rol Mobil", "Offline, bildirim, tema"),
        ("\U0001F310", "Multi-tenant", "Coklu kurum, izole veri"),
    ]

    cw = Inches(2.2)
    ch = Inches(1.55)
    gap = Inches(0.2)
    sx = Inches(0.55)
    sy = Inches(1.8)

    for i, (emoji, title, desc) in enumerate(features):
        row = i // 5
        col = i % 5
        x = sx + col * (cw + gap)
        y = sy + row * (ch + gap)

        colors = [BLUE, GREEN, PURPLE, GOLD, CYAN, ORANGE, PINK, TEAL, RED, EMERALD]
        accent = colors[i % len(colors)]

        _add_rounded_rect(slide, x, y, cw, ch,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1))
        _add_text_box(slide, x, y + Inches(0.1), cw, Inches(0.4),
                      emoji, Pt(22), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x + Inches(0.1), y + Inches(0.55),
                      cw - Inches(0.2), Inches(0.35),
                      title, Pt(11), WHITE, PP_ALIGN.CENTER, bold=True)
        _add_text_box(slide, x + Inches(0.1), y + Inches(0.9),
                      cw - Inches(0.2), Inches(0.5),
                      desc, Pt(9), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)

    # Bottom highlight
    _add_rounded_rect(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.8),
                      fill_color=CARD_BG_ALT, border_color=GOLD)
    _add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.8),
                  "\U0001F3AF  Bu 10 ozelligin HICBIRI rakip platformlarda bir arada bulunmaz",
                  Pt(14), GOLD, PP_ALIGN.CENTER, bold=True)


def slide_34_neden_biz():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "DEGER ONERISI")
    _slide_title(slide, "Neden Biz?", title_y=Inches(0.5))

    reasons = [
        ("\U0001F504", "TAM ENTEGRASYON",
         "30+ modul birbirine bagli calisir.\nVeri tek merkezde, tekrarsiz.\n"
         "Kayittan mezuniyete tek platform.\nHicbir veri kaybi veya tekrari yok.",
         BLUE),
        ("\U0001F916", "AI NATIVE",
         "Yapay zeka sonradan eklenmedi.\nTemelinden AI ile tasarlandi.\n"
         "Soru uretimi, risk analizi, chat.\nSurekli gelisen AI yetenekleri.",
         PURPLE),
        ("\U0001F3C6", "ENTERPRISE KALITE",
         "Dunya standartlarinda UI/UX.\nJWT + bcrypt guvenlik.\n"
         "Multi-tenant mimari.\n99.9% uptime hedefi.",
         GOLD),
    ]

    rw = Inches(3.6)
    rh = Inches(3.5)
    gap = Inches(0.35)
    sx = Inches(0.9)
    sy = Inches(1.8)

    for i, (emoji, title, desc, accent) in enumerate(reasons):
        x = sx + i * (rw + gap)
        _add_rounded_rect(slide, x, sy, rw, rh,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(2))
        # Number badge
        badge_size = Inches(0.5)
        _add_shape(slide, MSO_SHAPE.OVAL,
                   x + Inches(0.15), sy + Inches(0.15),
                   badge_size, badge_size,
                   fill_color=accent)
        _add_text_box(slide, x + Inches(0.15), sy + Inches(0.15),
                      badge_size, badge_size,
                      str(i + 1), Pt(16), WHITE, PP_ALIGN.CENTER, bold=True)
        # Emoji
        _add_text_box(slide, x + Inches(0.7), sy + Inches(0.2),
                      Inches(0.5), Inches(0.4),
                      emoji, Pt(28), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        # Title
        _add_text_box(slide, x + Inches(0.15), sy + Inches(0.8),
                      rw - Inches(0.3), Inches(0.35),
                      title, Pt(15), WHITE, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Gold line
        _add_line(slide, x + Inches(0.5), sy + Inches(1.2),
                  x + rw - Inches(0.5), sy + Inches(1.2),
                  color=accent, width=Pt(2))
        # Description
        _add_text_box(slide, x + Inches(0.25), sy + Inches(1.4),
                      rw - Inches(0.5), Inches(1.9),
                      desc, Pt(11), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False,
                      line_spacing=1.5)


# ====================================================================
# CLOSING SLIDES (35–40)
# ====================================================================

def slide_35_roi():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "YATIRIM GETIRISI")
    _slide_title(slide, "ROI & Deger Analizi", title_y=Inches(0.5))

    benefits = [
        ("\u23F0", "ZAMAN TASARRUFU",
         "Ogretmen basi haftada\n5+ saat tasarruf",
         "~%60", "Idari is yukundeki azalma", BLUE),
        ("\U0001F4C4", "KAGIT AZALMA",
         "Dijital sinav, yoklama,\nform ile kagit maliyeti",
         "~%80", "Kagit tuketimi azalma", GREEN),
        ("\U0001F60A", "VELI MEMNUNIYETI",
         "Aninda bilgilenme, randevu,\nodeme kolayligi",
         "~%40", "Memnuniyet artisi", PURPLE),
        ("\U0001F4C8", "KAYIT ARTISI",
         "Dijital imaj, hizli iletisim,\nprofesyonel algi",
         "~%25", "Yeni kayit artisi", GOLD),
    ]

    bw = Inches(2.7)
    bh = Inches(3.5)
    gap = Inches(0.3)
    sx = Inches(0.7)
    sy = Inches(1.7)

    for i, (emoji, title, desc, percent, metric, accent) in enumerate(benefits):
        x = sx + i * (bw + gap)
        _add_rounded_rect(slide, x, sy, bw, bh,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1.5))
        _add_text_box(slide, x, sy + Inches(0.15), bw, Inches(0.4),
                      emoji, Pt(28), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x, sy + Inches(0.55), bw, Inches(0.3),
                      title, Pt(12), WHITE, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        _add_text_box(slide, x + Inches(0.15), sy + Inches(0.95),
                      bw - Inches(0.3), Inches(0.7),
                      desc, Pt(10), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                      line_spacing=1.4)
        # Big percent
        _add_text_box(slide, x, sy + Inches(1.8), bw, Inches(0.7),
                      percent, Pt(40), accent, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Metric label
        _add_text_box(slide, x, sy + Inches(2.5), bw, Inches(0.35),
                      metric, Pt(10), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)

    # Bottom note
    _add_text_box(slide, Inches(0.5), Inches(5.7), Inches(12.0), Inches(0.5),
                  "* Tahminler sektordeki dijital donusum verilerine dayanmaktadir. Gercek sonuclar kuruma gore degisebilir.",
                  Pt(9), MED_GRAY, PP_ALIGN.CENTER, bold=False)


def slide_36_fiyat():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "FIYATLANDIRMA")
    _slide_title(slide, "Esnek Paketler", title_y=Inches(0.5))

    packages = [
        ("STARTER", "\u2B50", "Kucuk Okullar", [
            "10 temel modul",
            "1 yonetici + 10 ogretmen",
            "Mobil uygulama (2 rol)",
            "Email destek",
            "Aylik guncelleme",
        ], BLUE, "Aylik TBD"),
        ("PRO", "\U0001F48E", "Orta Olcekli Okullar", [
            "20 modul + AI ozellikler",
            "Sinirsiz kullanici",
            "Mobil uygulama (5 rol)",
            "SMS + Email entegrasyon",
            "Haftalik guncelleme",
            "Oncelikli destek",
        ], GOLD, "Aylik TBD"),
        ("ENTERPRISE", "\U0001F680", "Buyuk Egitim Kurumlari", [
            "30+ modul — TAM ERISIM",
            "Multi-tenant / coklu kampus",
            "Ozel entegrasyon",
            "7/24 destek + SLA",
            "Onsight egitim",
            "Kaynak kod lisansi opsiyonu",
        ], PURPLE, "Ozel Teklif"),
    ]

    pw = Inches(3.6)
    ph = Inches(4.5)
    gap = Inches(0.35)
    sx = Inches(0.9)
    sy = Inches(1.5)

    for i, (name, emoji, target, features, accent, price) in enumerate(packages):
        x = sx + i * (pw + gap)

        # Highlight PRO package
        is_pro = (i == 1)
        border_w = Pt(3) if is_pro else Pt(1.5)
        card_bg = CARD_BG_ALT if is_pro else CARD_BG

        _add_rounded_rect(slide, x, sy, pw, ph,
                          fill_color=card_bg, border_color=accent,
                          border_width=border_w)

        # Popular badge for PRO
        if is_pro:
            badge_w = Inches(1.5)
            _add_rounded_rect(slide, x + (pw - badge_w) / 2, sy - Inches(0.15),
                              badge_w, Inches(0.3),
                              fill_color=GOLD)
            _add_text_box(slide, x + (pw - badge_w) / 2, sy - Inches(0.15),
                          badge_w, Inches(0.3),
                          "EN POPULER", Pt(9), DARK_NAVY, PP_ALIGN.CENTER, bold=True)

        # Name
        _add_text_box(slide, x, sy + Inches(0.2), pw, Inches(0.4),
                      f"{emoji}  {name}", Pt(18), accent, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Target
        _add_text_box(slide, x, sy + Inches(0.65), pw, Inches(0.25),
                      target, Pt(11), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)

        # Divider
        _add_line(slide, x + Inches(0.3), sy + Inches(1.0),
                  x + pw - Inches(0.3), sy + Inches(1.0),
                  color=accent, width=Pt(1))

        # Features
        for j, feat in enumerate(features):
            _add_text_box(slide, x + Inches(0.2), sy + Inches(1.15) + j * Inches(0.38),
                          pw - Inches(0.4), Inches(0.35),
                          "\u2713  " + feat, Pt(10),
                          OFF_WHITE if not is_pro else GOLD,
                          PP_ALIGN.LEFT, bold=False)

        # Price
        _add_text_box(slide, x, sy + ph - Inches(0.6), pw, Inches(0.4),
                      price, Pt(14), accent, PP_ALIGN.CENTER, bold=True)

    # Bottom note
    _add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.0), Inches(0.3),
                  "Fiyatlar ogrenci sayisina ve secilen modullere gore belirlenir  |  Yillik odeme %15 indirim",
                  Pt(10), MED_GRAY, PP_ALIGN.CENTER, bold=False)


def slide_37_uygulama():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "UYGULAMA SURECI")
    _slide_title(slide, "4 Adimda Hayata Gecis", title_y=Inches(0.5))

    steps = [
        ("1", "DEMO", "\U0001F4BB",
         "Canli demo ile platformu\ntaniyin. Sorularinizi\nyonetim ekibimize sorun.",
         "1-2 Gun", BLUE),
        ("2", "KURULUM", "\u2699\uFE0F",
         "Kurum bilgileri, kullanici\ntanimlama, veri aktarimi.\nOzel yapilandirma.",
         "1-2 Hafta", GREEN),
        ("3", "EGITIM", "\U0001F393",
         "Yonetici, ogretmen ve\npersonel egitimi.\nKullanim kilavuzlari.",
         "1 Hafta", PURPLE),
        ("4", "GO-LIVE", "\U0001F680",
         "Canli ortama gecis.\n7/24 destek.\nIlk ay yogun takip.",
         "Surekli", GOLD),
    ]

    sw = Inches(2.5)
    sh = Inches(3.5)
    gap = Inches(0.5)
    sx = Inches(0.9)
    sy = Inches(1.7)

    for i, (num, title, emoji, desc, timeline, accent) in enumerate(steps):
        x = sx + i * (sw + gap)
        _add_rounded_rect(slide, x, sy, sw, sh,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1.5))
        # Number circle
        circle_size = Inches(0.55)
        _add_shape(slide, MSO_SHAPE.OVAL,
                   x + (sw - circle_size) / 2, sy + Inches(0.2),
                   circle_size, circle_size,
                   fill_color=accent)
        _add_text_box(slide, x + (sw - circle_size) / 2, sy + Inches(0.2),
                      circle_size, circle_size,
                      num, Pt(20), WHITE, PP_ALIGN.CENTER, bold=True)
        # Emoji
        _add_text_box(slide, x, sy + Inches(0.85), sw, Inches(0.4),
                      emoji, Pt(28), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        # Title
        _add_text_box(slide, x, sy + Inches(1.3), sw, Inches(0.3),
                      title, Pt(14), WHITE, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Description
        _add_text_box(slide, x + Inches(0.15), sy + Inches(1.7),
                      sw - Inches(0.3), Inches(1.0),
                      desc, Pt(10), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                      line_spacing=1.4)
        # Timeline badge
        badge_w = Inches(1.4)
        _add_rounded_rect(slide, x + (sw - badge_w) / 2, sy + sh - Inches(0.5),
                          badge_w, Inches(0.3),
                          fill_color=accent)
        _add_text_box(slide, x + (sw - badge_w) / 2, sy + sh - Inches(0.5),
                      badge_w, Inches(0.3),
                      timeline, Pt(9), WHITE, PP_ALIGN.CENTER, bold=True)

        # Arrow connector
        if i < 3:
            _add_text_box(slide, x + sw + Inches(0.05), sy + Inches(1.5),
                          Inches(0.4), Inches(0.4),
                          "\u25B6", Pt(20), GOLD, PP_ALIGN.CENTER, bold=False)

    # Bottom guarantee
    _add_rounded_rect(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.7),
                      fill_color=CARD_BG_ALT, border_color=GOLD)
    _add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.7),
                  "\U0001F91D  Uygulama sureci boyunca ozel proje yoneticisi atanir  |  7/24 teknik destek",
                  Pt(13), GOLD, PP_ALIGN.CENTER, bold=True)


def slide_38_referans():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "REFERANSLAR")
    _slide_title(slide, "Referanslar & Basari Hikayeleri", title_y=Inches(0.5))

    # Placeholder cards
    refs = [
        ("\U0001F3EB", "Ozel Koleji A", "Istanbul",
         "\"SmartCampus ile idari islerimiz %60 azaldi.\nVelilerimiz cok memnun.\"",
         "— Okul Muduru"),
        ("\U0001F3EB", "Ozel Koleji B", "Ankara",
         "\"AI soru uretimi ogretmenlerimizin\nis yukunu inanilmaz azaltti.\"",
         "— Akademik Koordinator"),
        ("\U0001F3EB", "Ozel Koleji C", "Izmir",
         "\"Mobil uygulama ile velilerimizle\niletisimimiz mukemmel seviyede.\"",
         "— Genel Mudur"),
    ]

    rw = Inches(3.6)
    rh = Inches(3.2)
    gap = Inches(0.35)
    sx = Inches(0.9)
    sy = Inches(1.8)

    for i, (emoji, name, city, quote, author) in enumerate(refs):
        x = sx + i * (rw + gap)
        _add_rounded_rect(slide, x, sy, rw, rh,
                          fill_color=CARD_BG, border_color=DARK_GOLD,
                          border_width=Pt(1))
        _add_text_box(slide, x, sy + Inches(0.2), rw, Inches(0.4),
                      emoji, Pt(32), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x, sy + Inches(0.65), rw, Inches(0.3),
                      name, Pt(14), WHITE, PP_ALIGN.CENTER, bold=True)
        _add_text_box(slide, x, sy + Inches(0.95), rw, Inches(0.25),
                      city, Pt(10), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)
        # Gold divider
        _add_line(slide, x + Inches(0.5), sy + Inches(1.3),
                  x + rw - Inches(0.5), sy + Inches(1.3),
                  color=GOLD, width=Pt(1))
        # Quote
        _add_text_box(slide, x + Inches(0.2), sy + Inches(1.5),
                      rw - Inches(0.4), Inches(1.0),
                      quote, Pt(10), OFF_WHITE, PP_ALIGN.CENTER, bold=False,
                      line_spacing=1.5, font_name="Segoe UI Light")
        # Author
        _add_text_box(slide, x + Inches(0.2), sy + Inches(2.6),
                      rw - Inches(0.4), Inches(0.3),
                      author, Pt(9), GOLD, PP_ALIGN.CENTER, bold=True)

    # Placeholder note
    _add_rounded_rect(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
                      fill_color=CARD_BG_ALT, border_color=MED_GRAY)
    _add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
                  "\U0001F4CB  Referans bilgileri guncellenmektedir.\n"
                  "Canli demo sirasinda mevcut kullanicilarin deneyimlerini paylasmaktayiz.",
                  Pt(11), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False, line_spacing=1.4)


def slide_39_iletisim():
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "ILETISIM")
    _slide_title(slide, "Bize Ulasin", title_y=Inches(0.5))

    # Contact card
    _add_rounded_rect(slide, Inches(2.0), Inches(1.8),
                      Inches(9.3), Inches(4.2),
                      fill_color=CARD_BG, border_color=GOLD,
                      border_width=Pt(2))

    # SmartCampus logo text
    _add_text_box(slide, Inches(2.0), Inches(2.1), Inches(9.3), Inches(0.6),
                  "\U0001F3EB  SmartCampus AI", Pt(28), GOLD, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Black")

    # Gold divider
    _add_line(slide, Inches(4.0), Inches(2.85), Inches(9.3), Inches(2.85),
              color=GOLD, width=Pt(2))

    contacts = [
        ("\U0001F4E7", "Email:", "info@smartcampus.com.tr"),
        ("\U0001F4F1", "Telefon:", "+90 (XXX) XXX XX XX"),
        ("\U0001F310", "Web:", "www.smartcampus.com.tr"),
        ("\U0001F4CD", "Adres:", "Istanbul, Turkiye"),
    ]

    cy = Inches(3.1)
    for emoji, label, value in contacts:
        _add_rich_text_box(
            slide, Inches(3.5), cy, Inches(6.0), Inches(0.4),
            [
                (emoji + "  ", Pt(18), GOLD, False, "Segoe UI Emoji"),
                (label + "  ", Pt(14), LIGHT_GRAY, True),
                (value, Pt(14), WHITE, False),
            ],
            alignment=PP_ALIGN.LEFT
        )
        cy += Inches(0.55)

    # Social / CTA
    _add_text_box(slide, Inches(2.0), Inches(5.3), Inches(9.3), Inches(0.4),
                  "Demo talep edin  |  Ucretsiz deneme  |  Ozel teklif alin",
                  Pt(13), GOLD, PP_ALIGN.CENTER, bold=True)

    # Bottom note
    _add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11.0), Inches(0.3),
                  "Satis ekibimiz 24 saat icerisinde size donecektir",
                  Pt(11), MED_GRAY, PP_ALIGN.CENTER, bold=False)


def slide_40_arka_kapak():
    slide = new_slide()
    add_gradient_bg(slide, DARK_NAVY, NAVY)

    # Top-left gold accent
    _add_line(slide, Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.5),
              color=GOLD, width=Pt(4))
    _add_line(slide, Inches(0.5), Inches(0.5), Inches(0.5), Inches(2.5),
              color=GOLD, width=Pt(4))

    # Main title
    _add_text_box(slide, Inches(1.0), Inches(2.2), Inches(11.0), Inches(1.0),
                  "SmartCampus AI", Pt(56), WHITE, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Black")

    # Gold divider
    _add_line(slide, Inches(4.5), Inches(3.4), Inches(8.833), Inches(3.4),
              color=GOLD, width=Pt(4))

    # Tagline
    _add_text_box(slide, Inches(1.0), Inches(3.7), Inches(11.0), Inches(0.6),
                  "Egitimin Gelecegi Burada", Pt(26), GOLD, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Light")

    # Subtitle
    _add_text_box(slide, Inches(2.0), Inches(4.5), Inches(9.0), Inches(0.5),
                  "Yatiriminiz icin tesekkur ederiz", Pt(16), LIGHT_GRAY,
                  PP_ALIGN.CENTER, bold=False)

    # Bottom-right gold accent
    _add_line(slide, Inches(9.833), Inches(6.6), Inches(12.833), Inches(6.6),
              color=GOLD, width=Pt(4))
    _add_line(slide, Inches(12.833), Inches(4.6), Inches(12.833), Inches(6.6),
              color=GOLD, width=Pt(4))

    # Copyright
    _add_text_box(slide, Inches(1.0), Inches(5.5), Inches(11.0), Inches(0.3),
                  "\u00A9 2026 SmartCampus AI  |  Tum haklari saklidir",
                  Pt(10), MED_GRAY, PP_ALIGN.CENTER, bold=False)


# ====================================================================
# BONUS SLIDES (41–50): Deep-dive extras
# ====================================================================

def slide_41_veri_modeli():
    """Data architecture overview."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: VERI MIMARISI")
    _slide_title(slide, "Veri Modeli & Depolama", title_y=Inches(0.5))

    models = [
        ("Akademik", "9 model", "students, grades, attendance, schedule..."),
        ("Olcme", "12 model", "questions, exams, sessions, results..."),
        ("Rehberlik", "6 model", "cases, interviews, mood, BEP..."),
        ("IK", "8 model", "employees, positions, leaves, interviews..."),
        ("Kayit", "5 model", "applications, leads, campaigns..."),
        ("Finans", "6 model", "payments, invoices, plans..."),
        ("Kurum", "7 model", "schedule, menu, transport, requests..."),
        ("Iletisim", "4 model", "messages, SMS logs, templates..."),
    ]

    cw = Inches(2.7)
    ch = Inches(1.1)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    sx = Inches(0.7)
    sy = Inches(1.7)

    colors = [BLUE, PURPLE, GREEN, ORANGE, CYAN, GOLD, TEAL, PINK]

    for i, (domain, count, desc) in enumerate(models):
        row = i // 4
        col = i % 4
        x = sx + col * (cw + gap_x)
        y = sy + row * (ch + gap_y)
        accent = colors[i]
        _add_rounded_rect(slide, x, y, cw, ch,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1))
        _add_rich_text_box(
            slide, x + Inches(0.15), y + Inches(0.1),
            cw - Inches(0.3), Inches(0.35),
            [
                (domain, Pt(12), WHITE, True),
                ("  |  ", Pt(9), MED_GRAY, False),
                (count, Pt(11), accent, True),
            ]
        )
        _add_text_box(slide, x + Inches(0.15), y + Inches(0.5),
                      cw - Inches(0.3), Inches(0.5),
                      desc, Pt(9), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)

    # Bottom: data flow
    _add_rounded_rect(slide, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.2),
                      fill_color=CARD_BG_ALT, border_color=DARK_GOLD)
    _add_text_box(slide, Inches(0.9), Inches(4.45),
                  Inches(11.5), Inches(0.3),
                  "VERI AKIS DIYAGRAMI", Pt(12), GOLD, PP_ALIGN.CENTER, bold=True)
    flow_text = (
        "Kullanici (Web/Mobil)  \u2192  API Gateway (FastAPI)  \u2192  Service Layer  \u2192  Data Store (JSON)\n"
        "              \u2193                          \u2193                       \u2193                    \u2193\n"
        "        Auth (JWT)            Validation           Business Logic      File System\n"
        "              \u2193                          \u2193                       \u2193                    \u2193\n"
        "        Session              Rate Limit             AI Services        Backup"
    )
    _add_text_box(slide, Inches(1.0), Inches(4.85),
                  Inches(11.3), Inches(1.5),
                  flow_text, Pt(10), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                  line_spacing=1.3, font_name="Consolas")


def slide_42_meb_uyum():
    """MEB compliance."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: MEB UYUMU")
    _slide_title(slide, "MEB Mevzuat Uyumu", title_y=Inches(0.5))

    items = [
        ("\u2705", "MEB Mufredat: 6839 yillik plan kaydi, 14 ders, 1-12. sinif"),
        ("\u2705", "Kazanim bazli takip: her ders icin MEB kazanim listesi"),
        ("\u2705", "e-Okul entegrasyonu: veri export/import uyumlulugu"),
        ("\u2705", "Sinav standartlari: LGS/TYT/AYT bolum yapilari"),
        ("\u2705", "Devamsizlik kurallari: MEB mevzuati ile uyumlu limitter"),
        ("\u2705", "Not sistemi: 100'luk + 5'lik + harf notu donusum"),
        ("\u2705", "Karne formati: MEB standart karne sablonu"),
        ("\u2705", "Disiplin sureci: MEB disiplin yonetmeligi adimlari"),
        ("\u2705", "Ozel okul yonetmeligi: idari ve mali surecler"),
        ("\u2705", "KVKK uyumlulugu: kisisel veri koruma gereksinimleri"),
    ]

    _add_icon_bullet_list(slide, Inches(0.8), Inches(1.7),
                          Inches(7.0), Inches(5.0), items,
                          font_size=Pt(12), color=OFF_WHITE, spacing=1.35)

    # Right side highlight
    _add_rounded_rect(slide, Inches(8.2), Inches(1.7),
                      Inches(4.3), Inches(4.5),
                      fill_color=CARD_BG_ALT, border_color=GOLD)
    _add_text_box(slide, Inches(8.4), Inches(1.85),
                  Inches(3.9), Inches(0.3),
                  "MEB VERI BANKASI", Pt(13), GOLD, PP_ALIGN.CENTER, bold=True)
    meb_stats = [
        ("6,839", "Yillik Plan Kaydi"),
        ("14", "Ders"),
        ("1-12", "Sinif Kademesi"),
        ("7", "Soru Tipi"),
        ("6", "Bloom Seviyesi"),
    ]
    stat_y = Inches(2.3)
    for num, label in meb_stats:
        _add_rich_text_box(
            slide, Inches(8.5), stat_y, Inches(3.7), Inches(0.4),
            [
                (num, Pt(18), GOLD, True),
                ("  " + label, Pt(11), LIGHT_GRAY, False),
            ],
            alignment=PP_ALIGN.CENTER
        )
        stat_y += Inches(0.5)


def slide_43_telafi():
    """Remediation system."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: TELAFI SISTEMI")
    _slide_title(slide, "Otomatik Telafi & Pekistirme", title_y=Inches(0.5))

    bands = [
        ("\U0001F534", "RED (0-49)", "Ozet + Quiz-1 (hemen) + Quiz-2 (48 saat)",
         "Ciddi eksiklik — yogun telafi", RED),
        ("\U0001F7E1", "YELLOW (50-69)", "Pekistirme (5 soru)",
         "Orta eksiklik — destekleyici", GOLD),
        ("\U0001F7E2", "GREEN (70-84)", "Haftalik tekrar (8 soru)",
         "Iyi seviye — surdurme", GREEN),
        ("\U0001F535", "BLUE (85-100)", "Zor set (5) + Hiz calismasi (10)",
         "Ustun — meydan okuma", BLUE),
    ]

    bw = Inches(11.5)
    bh = Inches(1.0)
    sx = Inches(0.9)
    sy = Inches(1.7)

    for i, (emoji, band, action, desc, accent) in enumerate(bands):
        y = sy + i * (bh + Inches(0.15))
        _add_rounded_rect(slide, sx, y, bw, bh,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1.5))
        _add_text_box(slide, sx + Inches(0.15), y + Inches(0.1),
                      Inches(0.4), Inches(0.4),
                      emoji, Pt(22), WHITE, PP_ALIGN.LEFT, bold=False,
                      font_name="Segoe UI Emoji")
        _add_rich_text_box(
            slide, sx + Inches(0.6), y + Inches(0.1),
            Inches(3.0), Inches(0.35),
            [
                (band, Pt(14), accent, True),
            ]
        )
        _add_text_box(slide, sx + Inches(0.6), y + Inches(0.5),
                      Inches(4.0), Inches(0.35),
                      desc, Pt(10), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)
        _add_text_box(slide, sx + Inches(5.0), y + Inches(0.15),
                      Inches(6.2), Inches(0.7),
                      action, Pt(11), OFF_WHITE, PP_ALIGN.LEFT, bold=False)

    # Bottom
    _add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.35),
                  "AutoGrader sinav sonrasi otomatik telafi gorevi olusturur  |  RemediationEngine ile entegre",
                  Pt(11), GOLD, PP_ALIGN.CENTER, bold=True)


def slide_44_sinav_tipleri():
    """Exam types."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: SINAV TIPLERI")
    _slide_title(slide, "7 Soru Tipi & Sinav Cesitleri", title_y=Inches(0.5))

    q_types = [
        ("\U0001F1E6", "MCQ", "Coktan secmeli (4-5 sik)"),
        ("\u2705", "True/False", "Dogru-Yanlis"),
        ("\u270D\uFE0F", "Fill Blank", "Bosluk doldurma"),
        ("\U0001F504", "Matching", "Eslestirme"),
        ("\U0001F522", "Ordering", "Siralama"),
        ("\U0001F4DD", "Cloze", "Metin ici bosluk"),
        ("\U0001F4D0", "Math Expr", "Matematik ifade"),
    ]

    qw = Inches(1.55)
    qh = Inches(1.5)
    gap = Inches(0.15)
    sx = Inches(0.5)
    sy = Inches(1.7)

    for i, (emoji, name, desc) in enumerate(q_types):
        x = sx + i * (qw + gap)
        _add_rounded_rect(slide, x, sy, qw, qh,
                          fill_color=CARD_BG, border_color=GOLD, border_width=Pt(1))
        _add_text_box(slide, x, sy + Inches(0.1), qw, Inches(0.35),
                      emoji, Pt(22), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x, sy + Inches(0.5), qw, Inches(0.3),
                      name, Pt(11), WHITE, PP_ALIGN.CENTER, bold=True)
        _add_text_box(slide, x + Inches(0.05), sy + Inches(0.85),
                      qw - Inches(0.1), Inches(0.5),
                      desc, Pt(8), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False)

    # Exam types below
    exam_types = [
        ("\U0001F4CB", "Yazili Sinav", "Klasik kagit sinav, PDF cikti, optik form"),
        ("\U0001F4BB", "Online Sinav", "Tab guvenlik, heartbeat, zamanlayici"),
        ("\U0001F916", "AI Uretim", "Kazanim bazli otomatik soru uretimi"),
        ("\U0001F4CA", "LGS/TYT/AYT", "OSYM formati, bolum yapilari"),
    ]

    ew = Inches(2.7)
    eh = Inches(1.6)
    egap = Inches(0.3)
    esx = Inches(0.7)
    esy = Inches(3.6)

    for i, (emoji, name, desc) in enumerate(exam_types):
        x = esx + i * (ew + egap)
        _feature_card(slide, x, esy, ew, eh, emoji, name, desc,
                      border_color=BLUE)


def slide_45_multi_tenant():
    """Multi-tenant architecture."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: MULTI-TENANT")
    _slide_title(slide, "Coklu Kurum Mimarisi", title_y=Inches(0.5))

    features = [
        ("\U0001F3EB", "Kurum Izolasyonu",
         "Her kurum kendi veri alaninda.\nBir kurumun verisi digerine\nasla gorulmez."),
        ("\U0001F512", "Guvenli Ayristirma",
         "Tenant ID bazli filtreleme.\nAPI seviyesinde zorunlu kontrol.\nVeri sizintisi imkansiz."),
        ("\U0001F527", "Kolay Ekleme",
         "Yeni kurum dakikalar icinde.\nSablon tabanli kurulum.\nOtomatik yapilandirma."),
        ("\U0001F4CA", "Merkezi Yonetim",
         "Super admin paneli.\nTum kurumlar tek ekrandan.\nToplu raporlama."),
    ]

    fw = Inches(2.7)
    fh = Inches(3.0)
    fgap = Inches(0.3)
    fsx = Inches(0.7)
    fsy = Inches(1.7)

    for i, (emoji, title, desc) in enumerate(features):
        x = fsx + i * (fw + fgap)
        colors = [BLUE, GREEN, PURPLE, GOLD]
        _feature_card(slide, x, fsy, fw, fh, emoji, title, desc,
                      border_color=colors[i])

    # Bottom diagram
    _add_rounded_rect(slide, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.3),
                      fill_color=CARD_BG_ALT, border_color=DARK_GOLD)
    _add_text_box(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.1),
                  "\U0001F3EB Kurum A  ←→  \U0001F5C4\uFE0F Veri A     |     "
                  "\U0001F3EB Kurum B  ←→  \U0001F5C4\uFE0F Veri B     |     "
                  "\U0001F3EB Kurum C  ←→  \U0001F5C4\uFE0F Veri C\n\n"
                  "                                       \u2B06\uFE0F  Merkezi API Gateway + Auth  \u2B06\uFE0F",
                  Pt(12), OFF_WHITE, PP_ALIGN.CENTER, bold=False, line_spacing=1.6)


def slide_46_gamification():
    """Gamification & student engagement."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: GAMIFICATION")
    _slide_title(slide, "Ogrenci Motivasyonu & Oyunlastirma", title_y=Inches(0.5))

    elements = [
        ("\U0001F3C6", "Rozetler", "Basari rozetleri ile\nbelirli hedeflere ulasma\nmotivasyonu", GOLD),
        ("\u2B50", "XP Puanlar", "Her aktivite icin\npuan kazanma,\nseviye atlama", BLUE),
        ("\U0001F4CA", "Lider Tablosu", "Haftalik/aylik\nsiralama listesi,\nsinif bazli", GREEN),
        ("\U0001F381", "Oduller", "Dijital ve fiziksel\noduller, kupon\nsistemi", PURPLE),
        ("\U0001F3AE", "Mini Oyunlar", "Matematik yarisi,\nkelime avı, hafiza,\nbulmaca", CYAN),
        ("\U0001F525", "Seri (Streak)", "Ardisik gun calisma,\nkayip korkusu,\nalişkanlık", ORANGE),
    ]

    ew = Inches(1.8)
    eh = Inches(2.8)
    egap = Inches(0.2)
    esx = Inches(0.55)
    esy = Inches(1.7)

    for i, (emoji, title, desc, accent) in enumerate(elements):
        x = esx + i * (ew + egap)
        _add_rounded_rect(slide, x, esy, ew, eh,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1))
        _add_text_box(slide, x, esy + Inches(0.15), ew, Inches(0.4),
                      emoji, Pt(28), WHITE, PP_ALIGN.CENTER, bold=False,
                      font_name="Segoe UI Emoji")
        _add_text_box(slide, x, esy + Inches(0.6), ew, Inches(0.3),
                      title, Pt(12), WHITE, PP_ALIGN.CENTER, bold=True)
        _add_text_box(slide, x + Inches(0.1), esy + Inches(1.0),
                      ew - Inches(0.2), Inches(1.5),
                      desc, Pt(9), LIGHT_GRAY, PP_ALIGN.CENTER, bold=False,
                      line_spacing=1.4)

    # Bottom insight
    _add_text_box(slide, Inches(0.5), Inches(5.0), Inches(12.0), Inches(0.5),
                  "Arastirmalar gamification ile ogrenci katiliminin %40-60 arttigini gostermektedir",
                  Pt(12), GOLD, PP_ALIGN.CENTER, bold=True)


def slide_47_proctoring():
    """Exam proctoring / security."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: SINAV GUVENLIGI")
    _slide_title(slide, "Online Sinav Gozetim Sistemi", title_y=Inches(0.5))

    security_features = [
        ("\U0001F4BB", "Browser Lockdown", "Tam ekran zorunlulugu.\nDiger uygulamalara gecis\nengellenir.", BLUE),
        ("\U0001F440", "Tab Izleme", "Tab degistirme tespit.\nOtomatik uyari.\nSinav gunlugu.", GREEN),
        ("\u23F0", "Heartbeat", "Periyodik baglanti kontrolu.\nKopma tespiti.\nOtomatik devam.", PURPLE),
        ("\U0001F6A8", "Risk Puanlama", "Supheli davranis skoru.\nOtomatik bayrak.\nInceleme kuyrugu.", RED),
    ]

    fw = Inches(2.7)
    fh = Inches(3.0)
    fgap = Inches(0.3)
    fsx = Inches(0.7)
    fsy = Inches(1.7)

    for i, (emoji, title, desc, accent) in enumerate(security_features):
        x = fsx + i * (fw + fgap)
        _feature_card(slide, x, fsy, fw, fh, emoji, title, desc,
                      border_color=accent)

    # Bottom stats
    stats_data = [
        ("Tab Switch", "Tespit & Kayit"),
        ("Heartbeat", "5 sn Aralik"),
        ("Risk Score", "0-100 Arasi"),
        ("Audit Log", "Tam Gecmis"),
    ]
    sx_stat = Inches(0.7)
    for num, label in stats_data:
        _add_small_stat(slide, sx_stat, Inches(5.2), Inches(2.7), Inches(0.8),
                        num, label, num_color=GOLD)
        sx_stat += Inches(3.0)


def slide_48_reporting():
    """Advanced reporting."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: RAPORLAMA")
    _slide_title(slide, "Kapsamli Raporlama Sistemi", title_y=Inches(0.5))

    report_categories = [
        ("AKADEMIK RAPORLAR", [
            "Sinif bazli basari raporu",
            "Ders bazli karsilastirma",
            "Ogrenci gelisim grafigi",
            "Karne & transkript",
        ], BLUE),
        ("IDARI RAPORLAR", [
            "Devamsizlik analizi",
            "Butce/finans raporu",
            "Personel performans",
            "Kayit istatistikleri",
        ], GREEN),
        ("AI RAPORLAR", [
            "Risk analiz raporu",
            "Tahmin modeli sonuclari",
            "Soru kalite raporu",
            "Ogrenme analitigi",
        ], PURPLE),
        ("YONETICI RAPORLAR", [
            "Dashboard ozet",
            "Donem karsilastirma",
            "KPI takip tablosu",
            "Trend analizi",
        ], GOLD),
    ]

    rw = Inches(2.7)
    rh = Inches(3.3)
    rgap = Inches(0.3)
    rsx = Inches(0.7)
    rsy = Inches(1.7)

    for i, (category, items, accent) in enumerate(report_categories):
        x = rsx + i * (rw + rgap)
        _add_rounded_rect(slide, x, rsy, rw, rh,
                          fill_color=CARD_BG, border_color=accent,
                          border_width=Pt(1.5))
        _add_text_box(slide, x, rsy + Inches(0.15), rw, Inches(0.3),
                      category, Pt(11), accent, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        _add_line(slide, x + Inches(0.3), rsy + Inches(0.55),
                  x + rw - Inches(0.3), rsy + Inches(0.55),
                  color=accent, width=Pt(1))
        for j, item in enumerate(items):
            _add_text_box(slide, x + Inches(0.15), rsy + Inches(0.7) + j * Inches(0.5),
                          rw - Inches(0.3), Inches(0.4),
                          "\u25C6  " + item, Pt(10), LIGHT_GRAY, PP_ALIGN.LEFT, bold=False)

    # Export formats
    _add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12.0), Inches(0.35),
                  "Export Formatlari:  PDF  |  Excel  |  CSV  |  Plotly Interaktif  |  Email Otomatik Gonderim",
                  Pt(12), GOLD, PP_ALIGN.CENTER, bold=True)


def slide_49_roadmap():
    """Product roadmap."""
    slide = new_slide()
    _section_label(slide, MARGIN_LEFT, Inches(0.3), "BONUS: YOLI HARITASI")
    _slide_title(slide, "Urun Yol Haritasi — 2026-2027", title_y=Inches(0.5))

    quarters = [
        ("Q2 2026", "MEVCUT", [
            "30+ web modul",
            "99 mobil sayfa",
            "143 API route",
            "AI soru uretimi",
        ], GOLD, True),
        ("Q3 2026", "YAKIN", [
            "iOS uygulamasi",
            "SQLite gecisi",
            "WhatsApp entegrasyon",
            "Video konferans",
        ], BLUE, False),
        ("Q4 2026", "PLANLANAN", [
            "Ogrenme analitigi v2",
            "Veli portal web",
            "E-tahsilat entegrasyon",
            "MEB e-Okul API",
        ], GREEN, False),
        ("Q1 2027", "VIZYON", [
            "Sesli AI asistan",
            "AR/VR ogrenme",
            "Blockchain sertifika",
            "Uluslararasi pazar",
        ], PURPLE, False),
    ]

    qw = Inches(2.7)
    qh = Inches(3.5)
    qgap = Inches(0.3)
    qsx = Inches(0.7)
    qsy = Inches(1.7)

    for i, (period, status, items, accent, is_current) in enumerate(quarters):
        x = qsx + i * (qw + qgap)
        border_w = Pt(3) if is_current else Pt(1)
        _add_rounded_rect(slide, x, qsy, qw, qh,
                          fill_color=CARD_BG if not is_current else CARD_BG_ALT,
                          border_color=accent, border_width=border_w)
        # Period
        _add_text_box(slide, x, qsy + Inches(0.15), qw, Inches(0.35),
                      period, Pt(16), accent, PP_ALIGN.CENTER, bold=True,
                      font_name="Segoe UI Black")
        # Status badge
        badge_w = Inches(1.3)
        _add_rounded_rect(slide, x + (qw - badge_w) / 2, qsy + Inches(0.55),
                          badge_w, Inches(0.25),
                          fill_color=accent)
        _add_text_box(slide, x + (qw - badge_w) / 2, qsy + Inches(0.55),
                      badge_w, Inches(0.25),
                      status, Pt(8), WHITE, PP_ALIGN.CENTER, bold=True)
        # Items
        for j, item in enumerate(items):
            marker = "\u2705" if is_current else "\u25CB"
            _add_text_box(slide, x + Inches(0.15), qsy + Inches(1.0) + j * Inches(0.5),
                          qw - Inches(0.3), Inches(0.4),
                          f"{marker}  {item}", Pt(10),
                          OFF_WHITE if is_current else LIGHT_GRAY,
                          PP_ALIGN.LEFT, bold=False)

        # Arrow
        if i < 3:
            _add_text_box(slide, x + qw + Inches(0.02), qsy + Inches(1.4),
                          Inches(0.3), Inches(0.3),
                          "\u25B6", Pt(16), GOLD, PP_ALIGN.CENTER, bold=False)


def slide_50_tesekkur():
    """Final thank-you slide."""
    slide = new_slide()
    add_gradient_bg(slide, DARK_NAVY, NAVY)

    # Gold frame accents (all 4 corners)
    # Top-left
    _add_line(slide, Inches(0.5), Inches(0.5), Inches(2.5), Inches(0.5),
              color=GOLD, width=Pt(3))
    _add_line(slide, Inches(0.5), Inches(0.5), Inches(0.5), Inches(2.0),
              color=GOLD, width=Pt(3))
    # Top-right
    _add_line(slide, Inches(10.833), Inches(0.5), Inches(12.833), Inches(0.5),
              color=GOLD, width=Pt(3))
    _add_line(slide, Inches(12.833), Inches(0.5), Inches(12.833), Inches(2.0),
              color=GOLD, width=Pt(3))
    # Bottom-left
    _add_line(slide, Inches(0.5), Inches(7.0), Inches(2.5), Inches(7.0),
              color=GOLD, width=Pt(3))
    _add_line(slide, Inches(0.5), Inches(5.5), Inches(0.5), Inches(7.0),
              color=GOLD, width=Pt(3))
    # Bottom-right
    _add_line(slide, Inches(10.833), Inches(7.0), Inches(12.833), Inches(7.0),
              color=GOLD, width=Pt(3))
    _add_line(slide, Inches(12.833), Inches(5.5), Inches(12.833), Inches(7.0),
              color=GOLD, width=Pt(3))

    # Diamond icon
    _add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.0), Inches(0.7),
                  "\U0001F48E", Pt(48), GOLD, PP_ALIGN.CENTER, bold=False,
                  font_name="Segoe UI Emoji")

    # Thank you
    _add_text_box(slide, Inches(1.0), Inches(2.3), Inches(11.0), Inches(0.8),
                  "Tesekkurler", Pt(48), WHITE, PP_ALIGN.CENTER, bold=True,
                  font_name="Segoe UI Black")

    # Gold divider
    _add_line(slide, Inches(5.0), Inches(3.3), Inches(8.333), Inches(3.3),
              color=GOLD, width=Pt(4))

    # Subtitle
    _add_text_box(slide, Inches(1.0), Inches(3.6), Inches(11.0), Inches(0.6),
                  "SmartCampus AI — Egitimin Gelecegi", Pt(24), GOLD,
                  PP_ALIGN.CENTER, bold=True, font_name="Segoe UI Light")

    # Call to action
    _add_text_box(slide, Inches(2.0), Inches(4.5), Inches(9.0), Inches(0.5),
                  "Sorulariniz icin hazirdayiz", Pt(16), LIGHT_GRAY,
                  PP_ALIGN.CENTER, bold=False)

    # Contact info
    _add_text_box(slide, Inches(2.0), Inches(5.3), Inches(9.0), Inches(0.4),
                  "\U0001F4E7 info@smartcampus.com.tr  |  \U0001F310 www.smartcampus.com.tr",
                  Pt(13), OFF_WHITE, PP_ALIGN.CENTER, bold=False)

    # Copyright
    _add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.0), Inches(0.3),
                  "\u00A9 2026 SmartCampus AI  |  Tum Hakları Saklidir  |  Gizli Dokuman",
                  Pt(9), MED_GRAY, PP_ALIGN.CENTER, bold=False)


# ====================================================================
# MASTER BUILD
# ====================================================================
def build_presentation():
    """Build all 50 slides and save."""
    print("\n" + "=" * 60)
    print("  SmartCampus AI — ULTRA PREMIUM Presentation Generator")
    print("=" * 60)

    slides = [
        ("01", "Kapak", slide_01_cover),
        ("02", "Neden SmartCampus?", slide_02_neden),
        ("03", "Cozum Ozeti", slide_03_cozum),
        ("04", "Rakamlarla SmartCampus", slide_04_numbers),
        ("05", "Roller", slide_05_roles),
        ("06", "Kayit Modulu", slide_06_kayit),
        ("07", "Akademik Takip", slide_07_akademik),
        ("08", "Olcme & Degerlendirme", slide_08_olcme),
        ("09", "Rehberlik", slide_09_rehberlik),
        ("10", "IK & Personel", slide_10_ik),
        ("11", "Butce & Odeme", slide_11_butce),
        ("12", "Kurum Hizmetleri", slide_12_kurum),
        ("13", "Erken Uyari", slide_13_uyari),
        ("14", "Veli Paneli", slide_14_veli),
        ("15", "Ogrenci Paneli", slide_15_ogrenci),
        ("16", "Dijital Ogrenme", slide_16_dijital),
        ("17", "Kutuphane & Etkinlik", slide_17_kutuphane),
        ("18", "Saglik & Guvenlik", slide_18_saglik),
        ("19", "Analitik Dashboard", slide_19_analitik),
        ("20", "Sosyal Medya & Iletisim", slide_20_sosyal),
        ("21", "Sertifika Uretici", slide_21_sertifika),
        ("22", "Destek Hizmetleri", slide_22_destek),
        ("23", "Toplanti & Randevu", slide_23_toplanti),
        ("24", "Mezunlar & Kariyer", slide_24_mezun),
        ("25", "AI Destek (SmartI)", slide_25_ai),
        ("26", "Mobil Uygulama", slide_26_mobil),
        ("27", "Mobil Ozellikler", slide_27_mobil_ozellik),
        ("28", "APK Build & CI/CD", slide_28_apk),
        ("29", "Teknik Altyapi", slide_29_teknik),
        ("30", "Guvenlik", slide_30_guvenlik),
        ("31", "Entegrasyonlar", slide_31_entegrasyon),
        ("32", "Rakip Karsilastirma", slide_32_rakip),
        ("33", "Fark Yaratan 10 Ozellik", slide_33_fark),
        ("34", "Neden Biz?", slide_34_neden_biz),
        ("35", "ROI & Deger", slide_35_roi),
        ("36", "Fiyatlandirma", slide_36_fiyat),
        ("37", "Uygulama Sureci", slide_37_uygulama),
        ("38", "Referanslar", slide_38_referans),
        ("39", "Iletisim", slide_39_iletisim),
        ("40", "Arka Kapak", slide_40_arka_kapak),
        ("41", "Veri Modeli", slide_41_veri_modeli),
        ("42", "MEB Uyumu", slide_42_meb_uyum),
        ("43", "Telafi Sistemi", slide_43_telafi),
        ("44", "Sinav Tipleri", slide_44_sinav_tipleri),
        ("45", "Multi-tenant", slide_45_multi_tenant),
        ("46", "Gamification", slide_46_gamification),
        ("47", "Sinav Guvenligi", slide_47_proctoring),
        ("48", "Raporlama", slide_48_reporting),
        ("49", "Yol Haritasi", slide_49_roadmap),
        ("50", "Tesekkurler", slide_50_tesekkur),
    ]

    for num, name, func in slides:
        print(f"  [{num}/50]  {name}...", end="")
        try:
            func()
            print("  OK")
        except Exception as e:
            print(f"  HATA: {e}")
            raise

    # Save
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "SmartCampusAI_Sunum_ULTRA.pptx")
    prs.save(out_path)
    print(f"\n{'=' * 60}")
    print(f"  BASARILI! {len(slides)} slayt olusturuldu.")
    print(f"  Dosya: {out_path}")
    print(f"  Boyut: {os.path.getsize(out_path) / 1024:.0f} KB")
    print(f"{'=' * 60}\n")
    return out_path


# ====================================================================
# ENTRY POINT
# ====================================================================
if __name__ == "__main__":
    build_presentation()
